from fastapi import FastAPI, UploadFile, File, Form
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import HTMLResponse, StreamingResponse
from fastapi.staticfiles import StaticFiles
import os
import json
import time
from io import BytesIO
from PIL import Image, ImageDraw, ImageFont, ImageOps
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from google import genai
from google.genai import types
from dotenv import load_dotenv

# 👆 ここまで置き換えたら、すぐ下は「# 🌟 追加：進捗をテキストファイルに書き込む関数」に繋がります

# 🌟 追加：進捗をテキストファイルに書き込む関数
def set_progress(message: str):
    try:
        with open("progress.txt", "w", encoding="utf-8") as f:
            f.write(message)
    except Exception as e:
        print(f"進捗の書き込みに失敗: {e}")
    
# --- 環境設定とGeminiの準備 ---
load_dotenv()
gemini_api_key = os.getenv("GEMINI_API_KEY")
# APIキーがあればGeminiクライアントを準備
gemini_client = genai.Client(api_key=gemini_api_key) if gemini_api_key else None
# --- 🌟ここから：データベース接続設定（SQLAlchemy） ---
from sqlalchemy import create_engine, Column, Integer, String
from sqlalchemy.orm import declarative_base, sessionmaker

DATABASE_URL = os.getenv("DATABASE_URL")
engine = create_engine(DATABASE_URL)
SessionLocal = sessionmaker(autocommit=False, autoflush=False, bind=engine)
Base = declarative_base()

# 店舗テーブルの設計図
class Branch(Base):
    __tablename__ = 'branches'
    id = Column(Integer, primary_key=True, autoincrement=True)
    branch_key = Column(String(50), unique=True, nullable=False)
    full_name = Column(String(100), nullable=False)
    license = Column(String(100), nullable=False)
    address = Column(String(200), nullable=False)
    tel = Column(String(20), nullable=False)
    login_id = Column(String(50), nullable=False)
    password = Column(String(100), nullable=False)
# --- 🌟ここまで ---
# 🌟 ここに追加：テーマテーブルの設計図
class Theme(Base):
    __tablename__ = 'themes'
    id = Column(Integer, primary_key=True, autoincrement=True)
    theme_key = Column(String(50), unique=True, nullable=False)
    name = Column(String(100), nullable=False)
    bg_color = Column(String(50), nullable=False)
    text_color = Column(String(50), nullable=False)
    accent_color = Column(String(50), nullable=False)

# 🌟 ここに追加：ユーザーテーブルの設計図
class User(Base):
    __tablename__ = 'users'
    id = Column(Integer, primary_key=True, autoincrement=True)
    email = Column(String(100), unique=True, nullable=False)
    name = Column(String(50), nullable=False)
    password = Column(String(100), nullable=False)
    branch_key = Column(String(50), nullable=False)
    footer_design_num = Column(String(10), default="1")
# 🌟 ここに追加：帯（フッター）履歴テーブルの設計図
from sqlalchemy import DateTime
from datetime import datetime

class FooterHistory(Base):
    __tablename__ = 'footer_history'
    id = Column(Integer, primary_key=True, autoincrement=True)
    user_email = Column(String(100), nullable=False)
    image_filename = Column(String(200), nullable=False)
    created_at = Column(DateTime, default=datetime.now)

FONT_PATH = './NotoSansCJKjp-Bold.ttf'

app = FastAPI(title="らくらく販売図面 APIサーバー")
# （これより下はそのまま残します）
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# 🌟 ここを追加：静的ファイル（CSS, JS）を配信するための設定
app.mount("/static", StaticFiles(directory="static"), name="static")

# 🌟 追加：画面(JS)から2秒ごとに呼ばれる進捗確認用API
@app.get("/apb_progress")
async def get_apb_progress():
    try:
        if os.path.exists("progress.txt"):
            with open("progress.txt", "r", encoding="utf-8") as f:
                return {"message": f.read()}
        return {"message": "🚀 処理を準備中..."}
    except Exception:
        return {"message": "進捗を確認中..."}

# 🌟 ブラウザからのアイコン要求（favicon）エラーを無視する設定
from fastapi.responses import Response
@app.get("/favicon.ico", include_in_schema=False)
async def favicon():
    return Response(status_code=204) # 「アイコンは無いから探さなくていいよ」と返事する

# ① ファイルの一番上（importが並んでいるところ）に以下を追加します
from fastapi.responses import HTMLResponse

# ② ログイン画面のルート部分に response_class=HTMLResponse を付け足します
@app.get("/", response_class=HTMLResponse)
def read_root():
    with open("templates/login.html", "r", encoding="utf-8") as f:
        return f.read()
# 👆 ここまで 👆

# 🌟 新設：個人アカウントでのログインAPI
@app.post("/api/login")
def api_login(username: str = Form(...), password: str = Form(...)):
    db = SessionLocal()
    try:
        user = db.query(User).filter(User.email == username, User.password == password).first()
        if user:
            return {
                "status": "success",
                "email": user.email,
                "name": user.name,
                "branch_key": user.branch_key,
                "footer_design_num": user.footer_design_num
            }
        else:
            return {"status": "error", "message": "メールアドレスまたはパスワードが間違っています。"}
    finally:
        db.close()

import shutil
import uuid

# 🌟 帯画像を保存する専用フォルダの設定（なければ自動作成）
FOOTER_DIR = os.path.join(os.path.dirname(__file__), "static", "footers")
os.makedirs(FOOTER_DIR, exist_ok=True)

# 🌟 API 1：帯画像をアップロードして保存する
@app.post("/api/upload_footer")
def upload_footer(user_email: str = Form(...), file: UploadFile = File(...)):
    db = SessionLocal()
    try:
        # ファイル名がかぶらないようにランダムなIDを生成
        ext = os.path.splitext(file.filename)[1]
        if not ext: ext = ".png"
        unique_filename = f"{uuid.uuid4().hex}{ext}"
        save_path = os.path.join(FOOTER_DIR, unique_filename)

        # サーバーの static/footers フォルダに実体画像を保存
        with open(save_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)

        # データベースに「誰がどの画像を保存したか」を記録
        new_history = FooterHistory(user_email=user_email, image_filename=unique_filename)
        db.add(new_history)
        db.commit()

        return {"status": "success", "filename": unique_filename, "url": f"/static/footers/{unique_filename}"}
    except Exception as e:
        db.rollback()
        return {"status": "error", "message": str(e)}
    finally:
        db.close()

# 🌟 API 2：個人の帯画像履歴を呼び出す（新しい順に最大10件）
@app.get("/api/footer_history/{user_email}")
def get_footer_history(user_email: str):
    db = SessionLocal()
    try:
        histories = db.query(FooterHistory).filter(FooterHistory.user_email == user_email).order_by(FooterHistory.created_at.desc()).limit(10).all()
        history_list = [
            {"id": h.id, "filename": h.image_filename, "url": f"/static/footers/{h.image_filename}"}
            for h in histories
        ]
        return {"status": "success", "histories": history_list}
    finally:
        db.close()
        
# 🌟 2. 総合メニューポータル画面
@app.get("/menu", response_class=HTMLResponse)
def read_menu():
    html_path = os.path.join(os.path.dirname(__file__), "templates", "menu.html")
    with open(html_path, "r", encoding="utf-8") as f:
        return f.read()

# 🌟 3. APB（パンフレット作成）画面
@app.get("/apb", response_class=HTMLResponse)
def read_apb():
    html_path = os.path.join(os.path.dirname(__file__), "templates", "apb.html")
    with open(html_path, "r", encoding="utf-8") as f:
        return f.read()

# 🌟 4. 販売図面作成画面（旧index.html）
@app.get("/zumen", response_class=HTMLResponse)
def read_zumen():
    html_path = os.path.join(os.path.dirname(__file__), "templates", "zumen.html")
    with open(html_path, "r", encoding="utf-8") as f:
        return f.read()
# ==========================================
# 🎨 パワポ装飾・画像処理用の便利ツール群
# ==========================================

def get_fitting_font(draw_obj, text, initial_size, max_width):
    """画像（PIL）に文字を入れる際、枠にはみ出さないようにサイズを自動調整する機能"""
    size = initial_size
    try: 
        font = ImageFont.truetype(FONT_PATH, size)
    except: 
        return ImageFont.load_default()
    while draw_obj.textbbox((0, 0), text, font=font)[2] > max_width and size > 10:
        size -= 2
        font = ImageFont.truetype(FONT_PATH, size)
    return font

def add_placeholder_box(slide_obj, left, top, width, height, text, luxury_gold_rgb=None, show_icon=False, label_text="MAIN SLOT"):
    """パワポ内に、高級感のあるゴールド枠線付きのプレースホルダー（画像挿入枠）を作る機能"""
    tx_box = slide_obj.shapes.add_textbox(left, top, width, height)
    tx_box.fill.solid()
    tx_box.fill.fore_color.rgb = RGBColor(245, 245, 245) # 予備の薄いグレー
    
    if luxury_gold_rgb:
        tx_box.line.color.rgb = luxury_gold_rgb
    else:
        tx_box.line.color.rgb = RGBColor(185, 160, 110) # デフォルトゴールド
    tx_box.line.width = Pt(1)
    
    tf = tx_box.text_frame
    tf.word_wrap = True
    tf.clear()
    
    if show_icon:
        p_label = tf.add_paragraph()
        p_label.text = label_text
        p_label.font.size = Pt(9)
        p_label.font.name = "游ゴシック"
        p_label.font.color.rgb = RGBColor(150, 150, 150)
        p_label.alignment = PP_ALIGN.CENTER
        
        p_icon = tf.add_paragraph()
        p_icon.text = "🔍"
        p_icon.font.size = Pt(12)
        p_icon.alignment = PP_ALIGN.CENTER
        
    p_jp = tf.add_paragraph()
    clean_text = text.replace("【", "").replace("】", "").replace(" ", "")
    p_jp.text = clean_text
    p_jp.font.size = Pt(11)
    p_jp.font.name = "游明朝"
    p_jp.font.color.rgb = RGBColor(80, 80, 80)
    p_jp.alignment = PP_ALIGN.CENTER

def add_station_info_circle(slide_obj, x, y, radius, text, color=(245, 240, 225)):
    """表紙などに使われる、駅徒歩情報の「円形グラフィック」を作る機能"""
    circle = slide_obj.shapes.add_shape(MSO_SHAPE.OVAL, x, y, radius*2, radius*2)
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(20, 20, 20)
    circle.line.color.rgb = RGBColor(color[0], color[1], color[2]) 
    circle.line.width = Pt(1.5)
    
    tf = circle.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.add_paragraph()
    p.text = text
    
    # 🌟 文字数に応じてフォントサイズを自動で縮小させる
    if len(text) > 16:
        p.font.size = Pt(9)
    elif len(text) > 12:
        p.font.size = Pt(10)
    else:
        p.font.size = Pt(11)
        
    p.font.color.rgb = RGBColor(color[0], color[1], color[2]) 
    p.font.name = "游明朝"
    p.alignment = PP_ALIGN.CENTER
    tf.margin_bottom = tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.05)
def generate_with_retry(model_name, prompt_contents, generation_config=None):
    """
    Geminiサーバーが混雑して429（回数制限）や503（一時停止）のエラーが出た際に、
    自動的に少し待ってから最大11回まで粘り強く再試行する、システムの守り神となる関数です。
    """
    wait_times = [30, 60, 120, 180, 240, 300, 360, 420, 480, 540, 600]
    max_attempts = len(wait_times) + 1
    
    for attempt in range(max_attempts):
        try:
            time.sleep(2)
            if generation_config:
                res = gemini_client.models.generate_content(model=model_name, contents=prompt_contents, config=generation_config)
            else:
                res = gemini_client.models.generate_content(model=model_name, contents=prompt_contents)
            return res
        except Exception as inner_e:
            # サーバー制限エラーを検知した場合、待機してリトライ
            if ("503" in str(inner_e) or "429" in str(inner_e)) and attempt < len(wait_times):
                wait_sec = wait_times[attempt]
                print(f"⚠️ Geminiの制限に達しました。{wait_sec}秒待機して自動再試行します... (試行 {attempt + 1}/{len(wait_times)})")
                set_progress(f"⚠️ AIサーバー混雑中... {wait_sec}秒待機して自動再開します (試行 {attempt + 1}/{len(wait_times)})")
                time.sleep(wait_sec)
            else:
                raise inner_e
@app.post("/generate_apb")
def generate_apb(
    zumen_file: UploadFile = File(None),
    madori_file: UploadFile = File(None),
    empty_file: UploadFile = File(None),
    map_file: UploadFile = File(None),
    orientation: str = Form("portrait"),
    selected_pages: list[str] = Form([])
): 
    try: # 🌟 ここから下の中身は、スペース4個分右に下がります
        import base64
        print("🚀 【APB】本番用のパンフレット生成プログラムを起動しました...")
        set_progress("🚀 処理を開始しています...")  # 🌟 最初の進捗メモ

        # 🌟 画面プレビュー用の軽量画像を格納するリスト
        preview_list = []
            
        def image_to_base64(img_obj):
            """Imageオブジェクトを画面表示用に小さく縮小し、軽量なJPEGに圧縮して変換する安全関数"""
            try:
                preview_img = img_obj.copy()
                preview_img.thumbnail((300, 300)) # 横幅最大300pxに縮小
                buf = BytesIO()
                preview_img.save(buf, format="JPEG", quality=50) # 高圧縮JPEGで激軽化
                return base64.b64encode(buf.getvalue()).decode('utf-8')
            except Exception as e:
                print(f"プレビュー画像変換エラー: {e}")
                dummy = Image.new('RGB', (100, 100), color=(200, 200, 200))
                buf = BytesIO()
                dummy.save(buf, format="JPEG")
                return base64.b64encode(buf.getvalue()).decode('utf-8')
            
        # 1. アップロードされた各種ファイルのバイナリ（データ）読み込み
        zumen_bytes = zumen_file.file.read() if zumen_file else None
        madori_bytes = madori_file.file.read() if madori_file else None
        empty_bytes = empty_file.file.read() if empty_file else None
        map_bytes = map_file.file.read() if map_file else None
        
        # 2. 【AI処理】販売図面のOCR解析（テキスト抽出）
        pdf_text = ""
        if zumen_bytes and gemini_client:
            print("🔍 販売図面をAI（OCR）で読み取り中...")
            set_progress("🔍 (1/6) 販売図面をAI（OCR）で読み取り中...")
            mime_type = "application/pdf" if zumen_file.filename.lower().endswith(".pdf") else "image/jpeg"
            ocr_analysis = generate_with_retry(
                model_name='gemini-2.5-flash',
                prompt_contents=[
                    "図面に記載されているすべてのテキスト情報を正確に抽出してください。特に「物件名」「所在地」「交通」「価格」「面積」「周辺施設」の情報は絶対に漏れなく抽出してください。Imagen生成用の地名（市区町村）も特定してください。",
                    types.Part.from_bytes(data=zumen_bytes, mime_type=mime_type)
                ]
            )
            pdf_text = ocr_analysis.text if ocr_analysis else ""

        # 3. 【AI処理】ページ構成用JSONデータの生成（プロの不動産ライター風）
        extracted_info = {}
        if pdf_text and gemini_client:
            print("🧠 AIが物件情報、地域名、デザインテーマを詳細分析中...")
            set_progress("🧠 (2/6) AIが物件情報やテーマを詳細分析中...")
            ratio_text = "4:3（横長）" if orientation == "landscape" else "3:4（縦長）"
            
            prompt = f"""
            あなたはプロの不動産ライターです。提供された【補助データ】を隅々まで解析し、以下の情報をJSON形式（オブジェクト、または配列の最初の要素）で出力してください。
            必須キー: property_name_jp, property_name_en, city_town, station_info, price, price_jp, land_area, building_area, sub_copy, headline, sub_headline, main_text, access_info, life_info, company_name, license, address, tel
            
            【厳守事項：文字数制限とフォーマット】
            パンフレットのレイアウト崩れを防ぐため、以下のルールを絶対に守ってください。
            ・sub_copy（表紙のキャッチコピー）: 最大35文字以内。簡潔に。
            ・station_info（円形アイコン用駅情報）: 綺麗に収まるよう短い3行で出力（例：「西武池袋線\\nひばりヶ丘駅\\n徒歩2分」）。
            ・sub_headline（サブ見出し）: 40文字以内で簡潔に。
            ・main_text（本文）: 100〜130文字程度に要約してください（最大でも絶対に150文字を超えないこと）。
            
            出力は必ず純粋なJSON形式のみにしてください。
            比率: {ratio_text}
            【補助データ】: {pdf_text}
            """
            
            json_res = generate_with_retry(
                model_name='gemini-2.5-flash', 
                prompt_contents=prompt, 
                generation_config=types.GenerateContentConfig(response_mime_type="application/json")
            )
            
            if json_res:
                raw_json = json_res.text.strip()
                if raw_json.startswith("```" + "json"): 
                    raw_json = raw_json[7:-3].strip()
                elif raw_json.startswith("```"): 
                    raw_json = raw_json[3:-3].strip()
                
                try:
                    parsed_data = json.loads(raw_json)
                    extracted_info = parsed_data[0] if isinstance(parsed_data, list) and len(parsed_data) > 0 else parsed_data
                except:
                    extracted_info = {}

# 🌟【完全クリーン版】URLの崩れとインデントを100%修正した実地図自動取得の魔法
        address_for_map = None
        if zumen_bytes and gemini_client and (not map_bytes):
            print("🔍 パワポ地図用に販売図面から正確な住所を直接抽出中...")
            try:
                mime_type = "application/pdf" if zumen_file.filename.lower().endswith(".pdf") else "image/jpeg"
                address_analysis = gemini_client.models.generate_content(
                    model='gemini-2.5-flash',
                    contents=[
                        "この図面から『物件の正確な所在地（住所）』を特定し、その住所の文字列だけを出力してください。余計な解説や文字は一切含めず、住所のみ（例：東京都西東京市谷戸町3-28-16）にしてください。",
                        types.Part.from_bytes(data=zumen_bytes, mime_type=mime_type)
                    ]
                )
                if address_analysis and address_analysis.text:
                    address_for_map = address_analysis.text.strip()
                    print(f"📍 パワポ地図用に特定した住所: {address_for_map}")
            except Exception as addr_err:
                print(f"⚠️ 図面からの住所直接抽出に失敗しました: {addr_err}")

        # 取得した正確な住所をもとに周辺の「本物のリアルな日本地図」を自動取得
        if address_for_map and gemini_client and (not map_bytes) and ("万円" not in address_for_map) and (len(address_for_map) > 3):
            print(f"🌍 住所「{address_for_map}」から周辺のリアル地図画像を自動生成中...")
            try:
                import urllib.parse
                import requests
                
                # 🌟【大改革】AI(Gemini)の推測による位置ズレ（ハルシネーション）を完全に排除！
                # 国土地理院が提供する「公式の住所検索API」を使用して、番地レベルの正確な緯度経度を一撃で取得します。
                search_url = f"https://msearch.gsi.go.jp/address-search/AddressSearch?q={urllib.parse.quote(address_for_map)}"
                geo_res = requests.get(search_url, timeout=5)
                
                lat, lng = None, None
                if geo_res.status_code == 200:
                    geo_data = geo_res.json()
                    if len(geo_data) > 0:
                        # 検索結果の1件目から正確な座標を取得 (GeoJSON形式: coordinates: [lng, lat])
                        lng = geo_data[0]["geometry"]["coordinates"][0]
                        lat = geo_data[0]["geometry"]["coordinates"][1]
                        print(f"📍 国土地理院APIから正確な座標を取得成功: {geo_data[0]['properties']['title']}")

                if lat and lng and lat != 0.0 and lng != 0.0 and (20.0 < lat < 50.0) and (120.0 < lng < 150.0):
                    print(f"📌 緯度経度を特定（lat={lat}, lng={lng}）。Googleマップのタイルを取得します...")
                    import math
                    
                    zoom = 16
                    lat_rad = math.radians(lat)
                    n = 2.0 ** zoom
                    
                    # 🌟 目的地の全世界での正確なピクセル座標を算出
                    x_pixel_global = (lng + 180.0) / 360.0 * n * 256
                    y_pixel_global = (1.0 - math.log(math.tan(lat_rad) + (1.0 / math.cos(lat_rad))) / math.pi) / 2.0 * n * 256
                    
                    xtile = int(x_pixel_global // 256)
                    ytile = int(y_pixel_global // 256)
                    
                    merged_img = Image.new('RGB', (256 * 3, 256 * 3))
                    tile_success_count = 0
                    
                    # 通常のブラウザからのアクセスに完璧に偽装するヘッダー（社内ネット対策）
                    browser_headers = {
                        "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36",
                        "Accept": "image/avif,image/webp,image/apng,image/svg+xml,image/*,*/*;q=0.8",
                        "Accept-Language": "ja,en-US;q=0.9,en;q=0.8"
                    }
                    
                    for i in range(-1, 2):
                        for j in range(-1, 2):
                            x = xtile + i
                            y = ytile + j
                            tile_url = f"https://mt1.google.com/vt/lyrs=m&x={x}&y={y}&z={zoom}"
                            try:
                                tile_res = requests.get(tile_url, headers=browser_headers, timeout=5)
                                if tile_res.status_code == 200:
                                    t_img = Image.open(BytesIO(tile_res.content)).convert("RGB")
                                    tile_success_count += 1
                                else:
                                    t_img = Image.new('RGB', (256, 256), color=(250, 250, 250))
                            except:
                                t_img = Image.new('RGB', (256, 256), color=(250, 250, 250))
                            merged_img.paste(t_img, ((i + 1) * 256, (j + 1) * 256))
                    
                    if tile_success_count > 0:
                        # 結合した画像の中で目的地がどこにあるかピクセル単位で計算
                        target_x_in_merged = x_pixel_global - (xtile - 1) * 256
                        target_y_in_merged = y_pixel_global - (ytile - 1) * 256
                        
                        # 目的地が中心(300, 250)になるように切り抜き枠を計算
                        crop_left = int(target_x_in_merged - 300)
                        crop_top = int(target_y_in_merged - 250)
                        
                        # 画像の枠外に出ないように安全ガード
                        crop_left = max(0, min(crop_left, 768 - 600))
                        crop_top = max(0, min(crop_top, 768 - 500))
                        
                        auto_map_img = merged_img.crop((crop_left, crop_top, crop_left + 600, crop_top + 500))
                        
                        # 🌟 正確なピン描画
                        px = target_x_in_merged - crop_left
                        py = target_y_in_merged - crop_top
                        
                        draw_pin = ImageDraw.Draw(auto_map_img)
                        draw_pin.ellipse([(px - 5, py - 2), (px + 5, py + 2)], fill=(180, 180, 180))
                        draw_pin.polygon([(px - 7, py - 13), (px + 7, py - 13), (px, py)], fill=(234, 67, 53))
                        draw_pin.ellipse([(px - 10, py - 31), (px + 10, py - 11)], fill=(234, 67, 53))
                        draw_pin.ellipse([(px - 4, py - 25), (px + 4, py - 17)], fill=(255, 255, 255))
                        
                        map_buf = BytesIO()
                        auto_map_img.save(map_buf, format="PNG")
                        map_bytes = map_buf.getvalue()
                        print("✅ 本物の周辺実地図（ピン中央配置）の自動取得・結合に成功しました！")
                    else:
                        print("⚠️ 会社のネットワーク制限により地図タイルのダウンロードが遮断されました。")
                else:
                    print("⚠️ 住所から緯度経度を特定できませんでした。")
            except Exception as map_err:
                import traceback
                print(f"⚠️ 地図画像の自動取得中にエラーが発生しました:\n{traceback.format_exc()}")
        cover_images = []
        if "cover" in selected_pages and gemini_client:
            print("📸 表紙のデザイン案を3パターン（家族・リビング等）同時にImagenで生成中...")
            set_progress("📸 (3/6) 表紙のデザイン案をAIで3パターン生成中...")
            target_place = extracted_info.get("city_town", "Tokyo")
            
            prompts = [
                f"High-end luxury magazine photography.\nA happy Japanese family playing happily in a beautiful sunny park in {target_place}.\nBright natural daylight, soft sunlight filtering through trees. Modern and sophisticated casual style.\nNO text.",
                f"High-end luxury lifestyle photography.\nA sophisticated Japanese man relaxing and drinking coffee in a modern, spacious luxury living room of a detached house in {target_place}.\nPremium interior design, close-up shot. NO text.",
                f"High-end luxury lifestyle photography.\nA happy Japanese family playing in a stylish, high-ceiling living room of a luxury modern apartment in {target_place}.\nSophisticated interior design, bright natural light, warm family atmosphere. NO text."
            ]
            
            for idx, p in enumerate(prompts):
                for attempt in range(4):
                    try:
                        img_res = gemini_client.models.generate_images(
                            model='imagen-4.0-generate-001',
                            prompt=p,
                            config=types.GenerateImagesConfig(
                                number_of_images=1, 
                                aspect_ratio="4:3" if orientation == "landscape" else "3:4"
                            )
                        )
                        cover_images.append(Image.open(BytesIO(img_res.generated_images[0].image.image_bytes)))
                        time.sleep(3)
                        break
                    except Exception as img_e:
                        if "429" in str(img_e) and attempt < 3:
                            print(f"⏳ 画像生成制限（429）を検知。15秒後に自動試行します... (再試行 {attempt+1}/3)")
                            set_progress(f"⏳ AI画像生成が混雑中... 15秒後に自動再試行します (再試行 {attempt+1}/3)")
                            time.sleep(15)
                        else:
                            raise img_e

        # 5. 【PowerPoint組み立て】
        print("📊 高級感のあるPowerPointスライドを1枚ずつ緻密に組み立て中...")
        set_progress("📊 (6/6) パワポファイルを組み立て中（最終仕上げ）...")
        prs = Presentation()
        if orientation == "landscape":
            prs.slide_width, prs.slide_height = Inches(10), Inches(7.5)
        else:
            prs.slide_width, prs.slide_height = Inches(7.5), Inches(10)

        width, height = (1024, 768) if orientation == "landscape" else (768, 1024)

        # ────────────────────────────────────────────────────────
        # 🌟 ① 表紙（Cover）の書き出し
        # ────────────────────────────────────────────────────────
        if "cover" in selected_pages:
            for idx, c_img in enumerate(cover_images if cover_images else [None]):
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                
                if c_img:
                    img_io = BytesIO()
                    c_img.save(img_io, format='PNG')
                    img_io.seek(0)
                    slide.shapes.add_picture(img_io, 0, 0, width=prs.slide_width, height=prs.slide_height)
                    preview_list.append({"type": f"P.1 表紙 (案{idx+1})", "image": image_to_base64(c_img)})
                else:
                    dummy = Image.new('RGB', (width, height), color=(40, 40, 45))
                    preview_list.append({"type": "P.1 表紙", "image": image_to_base64(dummy)})
                
                sw, sh = prs.slide_width, prs.slide_height
                gold_color = RGBColor(215, 185, 140)
                
                top_bar_h = Inches(1.8) if orientation == "landscape" else Inches(2.2)
                top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, sw, top_bar_h)
                top_bar.fill.solid()
                top_bar.fill.fore_color.rgb = RGBColor(20, 20, 20)
                try: top_bar.fill.transparency = 0.4
                except: pass
                top_bar.line.fill.background()
                
                tx_addr = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(4.0), Inches(0.4))
                tx_addr.text_frame.paragraphs[0].text = extracted_info.get("city_town", "TOKYO")
                tx_addr.text_frame.paragraphs[0].font.size = Pt(13)
                tx_addr.text_frame.paragraphs[0].font.name = "游明朝"
                tx_addr.text_frame.paragraphs[0].font.color.rgb = RGBColor(230, 230, 230)
                
                prop_name = extracted_info.get("property_name_jp", "【物件名】")
                logo_font_size = Pt(36) if len(prop_name) > 15 else (Pt(48) if len(prop_name) > 10 else Pt(60))
                
                tx_logo = slide.shapes.add_textbox(Inches(0.35), Inches(0.6), sw - Inches(0.8), Inches(2.0))
                p_logo = tx_logo.text_frame.paragraphs[0]
                p_logo.text = prop_name
                p_logo.font.size = logo_font_size
                p_logo.font.bold = True
                p_logo.font.name = "游明朝"
                p_logo.font.color.rgb = gold_color

                bottom_bar_h = Inches(1.2) if orientation == "landscape" else Inches(1.5)
                bottom_bar_y = sh - bottom_bar_h
                bottom_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, bottom_bar_y, sw, bottom_bar_h)
                bottom_bar.fill.solid()
                bottom_bar.fill.fore_color.rgb = RGBColor(20, 20, 20)
                try: bottom_bar.fill.transparency = 0.4
                except: pass
                bottom_bar.line.fill.background()
                
                # 🌟 Y座標を少し上げて2行になっても収まるようにする
                tx_sub = slide.shapes.add_textbox(Inches(1.5), bottom_bar_y + Inches(0.2), sw - Inches(1.8), Inches(1.0))
                tf_sub = tx_sub.text_frame
                tf_sub.word_wrap = True # 🌟 キャッチコピーの折り返しを有効化
                p_sub = tf_sub.paragraphs[0]
                p_sub.text = extracted_info.get("sub_copy", "都市の洗練と、静謐なるプライベートを纏う新邸。")
                p_sub.font.size = Pt(13)
                p_sub.font.name = "游明朝"
                p_sub.font.color.rgb = RGBColor(240, 240, 240)
                p_sub.alignment = PP_ALIGN.CENTER
                
                station_text = extracted_info.get("station_info", "最寄駅\n徒歩分数").replace('\\n', '\n')
                # 🌟 円のサイズを少し大きく（0.55 → 0.65）し、位置を微調整して文字あふれを防ぐ
                add_station_info_circle(slide, Inches(0.35), bottom_bar_y - Inches(0.6), Inches(0.65), station_text, color=(215, 185, 140))

        # ────────────────────────────────────────────────────────
        # 🌟 ② コンセプト（Concept / aerial_map）の書き出し
        # ────────────────────────────────────────────────────────
        if "aerial_map" in selected_pages:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(245, 240, 225)
            
            headline = "THE CONCEPT"
            sub_headline = extracted_info.get('sub_headline', '五感を満たす、静謐の邸宅。')
            main_text = extracted_info.get('main_text', '')
            bronze_rgb = RGBColor(215, 185, 140)
            
            concept_bg = Image.new('RGB', (width, height), color=(245, 240, 225))
            preview_list.append({"type": "P.2 コンセプト", "image": image_to_base64(concept_bg)})
            
            if orientation == "landscape":
                tx_box_head = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(4.5), Inches(0.6))
                line1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.3), Inches(4.5), Pt(1))
                line1.fill.solid()
                line1.fill.fore_color.rgb = bronze_rgb
                line1.line.color.rgb = bronze_rgb
                tx_box_sub = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(4.5), Inches(0.6))
                tx_box_main = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(4.0), Inches(3.0))
                
                add_placeholder_box(slide, Inches(4.5), Inches(0), Inches(5.5), Inches(7.5), "【 メイン画像 】")
                add_placeholder_box(slide, Inches(3.0), Inches(4.5), Inches(2.5), Inches(2.0), "写真1")
                add_placeholder_box(slide, Inches(7.5), Inches(0.5), Inches(2.0), Inches(2.5), "写真2")
            else:
                tx_box_head = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(6.5), Inches(0.6))
                line1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.3), Inches(6.5), Pt(1))
                line1.fill.solid()
                line1.fill.fore_color.rgb = bronze_rgb
                line1.line.color.rgb = bronze_rgb
                tx_box_sub = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6.5), Inches(0.8)) # 🌟 高さを広げて2行対応
                tx_box_main = slide.shapes.add_textbox(Inches(0.5), Inches(2.4), Inches(6.5), Inches(2.0))
                
                add_placeholder_box(slide, Inches(0), Inches(4.5), Inches(7.5), Inches(5.5), "【 メイン画像 】")
                # 🌟 写真1のY座標(縦位置)を 3.5 → 4.8 に下げて、本文と被らないように修正
                add_placeholder_box(slide, Inches(4.5), Inches(4.8), Inches(2.5), Inches(2.0), "写真1")
                add_placeholder_box(slide, Inches(0.5), Inches(7.5), Inches(2.5), Inches(2.0), "写真2")

            p_head = tx_box_head.text_frame.paragraphs[0]
            p_head.text = headline
            p_head.font.name = "Arial"
            p_head.font.size = Pt(36)
            p_head.font.bold = True
            
            # 🌟 サブヘッドラインにも Word Wrap（折り返し）を追加し、サイズを少し落として安全マージンを取る
            tf_sub = tx_box_sub.text_frame
            tf_sub.word_wrap = True
            p_sub = tf_sub.paragraphs[0]
            p_sub.text = sub_headline
            p_sub.font.name = "游ゴシック"
            p_sub.font.size = Pt(20) # (元は24)
            p_sub.font.bold = True
            p_sub.font.color.rgb = RGBColor(0,0,0)

            tf_main = tx_box_main.text_frame
            tf_main.word_wrap = True
            p_main = tf_main.paragraphs[0]
            p_main.text = main_text
            p_main.font.name = "游ゴシック"
            p_main.font.size = Pt(12) # 🌟 文字が多い場合に備えて少し小さく (元は13)
            p_main.font.color.rgb = RGBColor(80,80,80)

        # ────────────────────────────────────────────────────────
        # 🌟 ③ MAP & ACCESS（アクセス・地図）の書き出し
        # ────────────────────────────────────────────────────────
        if "access" in selected_pages:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            acc_text = extracted_info.get('access_info', "交通情報がありません")
            if isinstance(acc_text, list): acc_text = "\n".join(str(x) for x in acc_text)
            
            life_text = extracted_info.get('life_info', "周辺施設情報がありません")
            if isinstance(life_text, list): life_text = "\n".join(str(x) for x in life_text)

            gold_color = RGBColor(180, 150, 80)
            navy_color = RGBColor(20, 30, 60)
            text_color_rgb = RGBColor(40, 40, 40)
            
            gen_map_io = None
            access_bg = Image.new('RGB', (width, height), color=(250, 250, 250))
            if map_bytes:
                print("🗺️ 地図画像をそのままのカラーで読み込み中...")
                try:
                    final_map = Image.open(BytesIO(map_bytes)).convert("RGB")
                    
                    gen_map_io = BytesIO()
                    final_map.save(gen_map_io, format='PNG')
                    gen_map_io.seek(0)
                    
                    resized_map = final_map.resize((int(width * 0.7), int(height * 0.5)))
                    access_bg.paste(resized_map, (int(width * 0.15), int(height * 0.2)))
                except Exception as e:
                    print(f"地図読み込みエラー: {e}")
            preview_list.append({"type": "P.3 地図", "image": image_to_base64(access_bg)})

            if orientation == "landscape":
                tx_box_head = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(5.0), Inches(0.8))
                tx_box_sub = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(5.0), Inches(0.4))
                line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(9.0), Pt(1.5))
                line.fill.solid()
                line.fill.fore_color.rgb = gold_color
                line.line.color.rgb = gold_color

                if gen_map_io:
                    slide.shapes.add_picture(gen_map_io, Inches(0.5), Inches(1.7), width=Inches(9.0), height=Inches(4.2))
                else:
                    add_placeholder_box(slide, Inches(0.5), Inches(1.7), Inches(9.0), Inches(4.2), "【 MAP画像 挿入枠 】\n※地図画像をアップロードすると自動で高級化されます")

                tx_box_acc = slide.shapes.add_textbox(Inches(0.5), Inches(6.1), Inches(4.3), Inches(1.2))
                tx_box_life = slide.shapes.add_textbox(Inches(5.2), Inches(6.1), Inches(4.3), Inches(1.2))
            else:
                tx_box_head = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(6.5), Inches(0.8))
                tx_box_sub = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(6.5), Inches(0.4))
                line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.6), Inches(6.5), Pt(1.5))
                line.fill.solid()
                line.fill.fore_color.rgb = gold_color
                line.line.color.rgb = gold_color

                if gen_map_io:
                    slide.shapes.add_picture(gen_map_io, Inches(0.5), Inches(1.8), width=Inches(6.5), height=Inches(5.0))
                else:
                    add_placeholder_box(slide, Inches(0.5), Inches(1.8), Inches(6.5), Inches(5.0), "【 MAP画像 挿入枠 】\n※地図画像をアップロードすると自動で高級化されます")
                
                tx_box_acc = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(3.1), Inches(2.5))
                tx_box_life = slide.shapes.add_textbox(Inches(3.9), Inches(7.0), Inches(3.1), Inches(2.5))

            p_head = tx_box_head.text_frame.paragraphs[0]
            p_head.text = "MAP & ACCESS"
            p_head.font.size = Pt(40)
            p_head.font.bold = True
            p_head.font.color.rgb = navy_color
            p_head.font.name = "Arial"

            p_sub = tx_box_sub.text_frame.paragraphs[0]
            p_sub.text = "周辺環境・アクセス"
            p_sub.font.size = Pt(20)
            p_sub.font.bold = True
            p_sub.font.color.rgb = gold_color
            p_sub.font.name = "游ゴシック"

            tf_acc = tx_box_acc.text_frame
            tf_acc.word_wrap = True
            tf_acc.clear() 
            p_acc_head = tf_acc.add_paragraph()
            p_acc_head.text = "■ 交通アクセス"
            p_acc_head.font.size = Pt(16)
            p_acc_head.font.bold = True
            p_acc_head.font.color.rgb = gold_color 
            p_acc_head.font.name = "游ゴシック"
            p_acc_body = tf_acc.add_paragraph()
            p_acc_body.text = acc_text
            p_acc_body.font.size = Pt(13)
            p_acc_body.font.color.rgb = text_color_rgb 
            p_acc_body.font.name = "游ゴシック"

            tf_life = tx_box_life.text_frame
            tf_life.word_wrap = True
            tf_life.clear() 
            p_life_head = tf_life.add_paragraph()
            p_life_head.text = "■ Life Information"
            p_life_head.font.size = Pt(16)
            p_life_head.font.bold = True
            p_life_head.font.color.rgb = gold_color 
            p_life_head.font.name = "Arial"
            p_life_body = tf_life.add_paragraph()
            p_life_body.text = life_text
            p_life_body.font.size = Pt(13)
            p_life_body.font.color.rgb = text_color_rgb
            p_life_body.font.name = "游ゴシック"

        # ────────────────────────────────────────────────────────
        # 🌟 ④ 間取り（Floor Plan）の書き出し（修正版）
        # ────────────────────────────────────────────────────────
        if "floor_plan" in selected_pages:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            luxury_gold = RGBColor(185, 160, 110)
            
            print("🎨 間取り図ページの高級背景（木目グラデーション）をAIで生成中...")
            set_progress("🎨 (4/6) 間取り図ページの背景画像を生成中...")
            bg_prompt = (
                "A cinematic, high-end architectural background for a luxury property brochure.\n"
                "Top 20% area: dark charcoal gray slate stone with elegant natural texture.\n"
                "Bottom 80% area: rich, deep espresso luxury wood grain texture, vertical planks, "
                "with a subtle vertical gradient that gets progressively darker towards the bottom edge.\n"
                "High contrast, moody lighting, NO text, NO floor plans, NO logos."
            )
            
            try:
                image_result = gemini_client.models.generate_images(
                    model='imagen-4.0-generate-001',
                    prompt=bg_prompt,
                    config=types.GenerateImagesConfig(
                        number_of_images=1, 
                        aspect_ratio="4:3" if orientation == "landscape" else "3:4"
                    )
                )
                generated_bytes = image_result.generated_images[0].image.image_bytes
                fp_bg = Image.open(BytesIO(generated_bytes)).convert("RGB").resize((width, height))
            except Exception as e:
                print(f"4ページ目の背景生成エラー: {e}")
                fp_bg = Image.new('RGB', (width, height), color=(25, 20, 15))

            draw = ImageDraw.Draw(fp_bg)
            header_h = int(height * 0.20)

            # ゴールドの区切り線を立体的に描画
            gold_light = (215, 185, 140)
            gold_dark = (160, 130, 80)
            draw.rectangle([(0, header_h - 15), (width, header_h + 15)], fill=gold_light)
            draw.rectangle([(0, header_h - 2), (width, header_h + 2)], fill=gold_dark)

            # 間取り図を立体ゴールド額縁付きで合成
            if madori_bytes:
                try:
                    m_img = Image.open(BytesIO(madori_bytes)).convert("RGB")
                    frame_width = 15
                    m_img_with_frame = Image.new('RGB', (m_img.width + frame_width*2, m_img.height + frame_width*2), gold_light)
                    draw_frame = ImageDraw.Draw(m_img_with_frame)
                    draw_frame.rectangle([(5, 5), (m_img_with_frame.width-5, m_img_with_frame.height-5)], outline=(100, 80, 40), width=3)
                    m_img_with_frame.paste(m_img, (frame_width, frame_width))
                    m_img_with_frame.thumbnail((width * 0.85, height * 0.70))
                    
                    p_x = (width - m_img_with_frame.width) // 2
                    p_y = header_h + 40 + (height - header_h - 40 - m_img_with_frame.height) // 2
                    fp_bg.paste(m_img_with_frame, (int(p_x), int(p_y)))
                except Exception as e:
                    print(f"間取り図合成エラー: {e}")
            else:
                try:
                    f_m = ImageFont.truetype(FONT_PATH, int(height * 0.025))
                    draw.text((width/2, height/2), "間取り図がアップロードされていません", font=f_m, fill="white", anchor="mm")
                except: pass

            preview_list.append({"type": "P.4 間取り図", "image": image_to_base64(fp_bg)})

            # 背景画像貼り付け
            img_io_fp = BytesIO()
            fp_bg.save(img_io_fp, format='PNG')
            img_io_fp.seek(0)
            slide.shapes.add_picture(img_io_fp, 0, 0, width=prs.slide_width, height=prs.slide_height)

            # 【テキストボックスの配置最適化】
            if orientation == "landscape":
                tb_title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(5.8), Inches(1.2))
                tb_area = slide.shapes.add_textbox(Inches(6.5), Inches(0.20), Inches(3.0), Inches(0.7))
                tb_price = slide.shapes.add_textbox(Inches(6.0), Inches(0.85), Inches(3.5), Inches(0.8))
            else:
                tb_title = slide.shapes.add_textbox(Inches(0.4), Inches(0.3), Inches(4.2), Inches(1.3))
                tb_area = slide.shapes.add_textbox(Inches(4.6), Inches(0.25), Inches(2.5), Inches(0.6))
                tb_price = slide.shapes.add_textbox(Inches(4.0), Inches(0.85), Inches(3.1), Inches(0.8))
                
            # ────────── 物件タイトル ──────────
            tf_title = tb_title.text_frame
            tf_title.word_wrap = True
            tf_title.margin_top = tf_title.margin_bottom = tf_title.margin_left = tf_title.margin_right = 0
            p_t = tf_title.paragraphs[0]
            p_t.text = extracted_info.get("property_name_jp", "物件名")
            
            name_len = len(p_t.text)
            if name_len > 20: p_t.font.size = Pt(20)
            elif name_len > 14: p_t.font.size = Pt(24)
            elif name_len > 10: p_t.font.size = Pt(28)
            else: p_t.font.size = Pt(32)
                
            p_t.font.name = "游明朝"
            p_t.font.bold = True
            p_t.font.color.rgb = luxury_gold
            try: p_t.font.shadow = True
            except: pass
            
            # ────────── 土地・建物面積（AIデータの自動クレンジング装置） ──────────
            import re
            tf_area = tb_area.text_frame
            tf_area.clear()
            tf_area.margin_top = tf_area.margin_bottom = tf_area.margin_left = tf_area.margin_right = 0
            
            p_l = tf_area.paragraphs[0]
            l_area = extracted_info.get("land_area")
            l_area_str = str(l_area).strip() if l_area is not None else ""
            
            if "所有権" in l_area_str or not re.search(r'\d', l_area_str):
                p_l.text = "土地 --- ㎡"
            else:
                if "土地" not in l_area_str: l_area_str = "土地 " + l_area_str
                if "㎡" not in l_area_str and "坪" not in l_area_str: l_area_str = l_area_str + "㎡"
                p_l.text = l_area_str
                
            p_l.font.size = Pt(14)
            p_l.font.color.rgb = luxury_gold
            p_l.alignment = PP_ALIGN.RIGHT
            try: p_l.font.shadow = True
            except: pass
            
            b_area = extracted_info.get("building_area")
            b_area_str = str(b_area).strip() if b_area is not None else ""
            
            if re.search(r'\d', b_area_str):
                p_b = tf_area.add_paragraph()
                if "建物" not in b_area_str: b_area_str = "建物 " + b_area_str
                if "㎡" not in b_area_str and "坪" not in b_area_str: b_area_str = b_area_str + "㎡"
                p_b.text = b_area_str
                p_b.font.size = Pt(14)
                p_b.font.color.rgb = luxury_gold
                p_b.alignment = PP_ALIGN.RIGHT
                try: p_b.font.shadow = True
                except: pass
                
            # ────────── 価格 ──────────
            tf_price = tb_price.text_frame
            tf_price.margin_top = tf_price.margin_bottom = tf_price.margin_left = tf_price.margin_right = 0
            p_p = tf_price.paragraphs[0]
            
            price_text = extracted_info.get("price_jp", "--- 万円")
            p_p.text = price_text
            p_p.font.name = "游明朝"
            p_p.font.bold = True
            p_p.font.color.rgb = luxury_gold
            p_p.alignment = PP_ALIGN.RIGHT
            try: p_p.font.shadow = True
            except: pass
            
            if len(price_text) > 12: p_p.font.size = Pt(24)
            elif len(price_text) > 8: p_p.font.size = Pt(30)
            else: p_p.font.size = Pt(36)

        # ────────────────────────────────────────────────────────
        # 🌟 ⑤ 内観生成（Interior HQ / バーチャルステージング）の書き出し
        # ────────────────────────────────────────────────────────
        if "interior_hq" in selected_pages:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            final_room_prompt = "A cinematic, high-end luxury living room interior, modern sophisticated design, bright natural light, luxury furniture, 8k resolution.\nNO text."
            if empty_bytes and gemini_client:
                print("🔍 空室写真を解析し、バーチャルステージング用のプロンプトを生成中...")
                try:
                    analysis_img = generate_with_retry(
                        model_name='gemini-2.5-flash',
                        prompt_contents=[
                            "Analyze this empty room photo precisely. Generate a highly detailed English prompt for an Image Generation AI to recreate this EXACT room structure (same window placement, floor, wall) BUT add elegant luxury style furniture. Output ONLY the prompt string.",
                            types.Part.from_bytes(data=empty_bytes, mime_type="image/jpeg")
                        ]
                    )
                    if analysis_img and analysis_img.text:
                        final_room_prompt = analysis_img.text.strip()
                except Exception as e:
                    print(f"空室解析エラー: {e}")

            bg_image = None
            if gemini_client:
                print("🎨 内観完成予想イメージをImagenで生成中...")
                set_progress("🎨 (5/6) 内観イメージをAI家具合成で生成中（約1分かかります）...")
                for attempt in range(3):
                    try:
                        image_result = gemini_client.models.generate_images(
                            model='imagen-4.0-generate-001',
                            prompt=final_room_prompt,
                            config=types.GenerateImagesConfig(
                                number_of_images=1, 
                                aspect_ratio="4:3" if orientation == "landscape" else "3:4"
                            )
                        )
                        generated_bytes = image_result.generated_images[0].image.image_bytes
                        bg_image = Image.open(BytesIO(generated_bytes)).convert("RGB")
                        break
                    except Exception as e:
                        if "429" in str(e) and attempt < 2:
                            print(f"⏳ 内観生成制限（429）を検知。15秒後に自動試行します... (再試行 {attempt+1}/2)")
                            set_progress(f"⏳ 内観画像の生成が混雑中... 15秒後に自動再試行します (再試行 {attempt+1}/2)")
                            time.sleep(15)
                        else:
                            print(f"内観生成エラー: {e}")
                            break
            
            if bg_image:
                img_io = BytesIO()
                bg_image.save(img_io, format='PNG')
                img_io.seek(0)
                slide.shapes.add_picture(img_io, 0, 0, width=prs.slide_width, height=prs.slide_height)
                preview_list.append({"type": "P.5 内観完成予想", "image": image_to_base64(bg_image)})
            else:
                dummy = Image.new('RGB', (width, height), color=(240, 240, 240))
                preview_list.append({"type": "P.5 内観完成予想", "image": image_to_base64(dummy)})

            sw, sh = prs.slide_width, prs.slide_height
            bar_h = Inches(1.5)
            bar_y = sh - bar_h
            bottom_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, bar_y, sw, bar_h)
            bottom_bar.fill.solid()
            bottom_bar.fill.fore_color.rgb = RGBColor(0, 0, 0)
            try: bottom_bar.fill.transparency = 0.3
            except: pass
            bottom_bar.line.fill.background()
            
            gold_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, bar_y, sw, Pt(4))
            gold_line.fill.solid()
            gold_line.fill.fore_color.rgb = RGBColor(185, 160, 110)
            gold_line.line.fill.background()
            
            tx_head = slide.shapes.add_textbox(sw * 0.05, bar_y + Inches(0.15), sw * 0.5, Inches(0.8))
            tf_head = tx_head.text_frame

            p_en = tf_head.paragraphs[0]
            p_en.text = "INTERIOR VISION"
            p_en.font.size = Pt(30)
            p_en.font.color.rgb = RGBColor(255, 255, 255)
            p_en.alignment = PP_ALIGN.LEFT

            p_jp = tf_head.add_paragraph()
            p_jp.text = "内観完成予想イメージ"
            p_jp.font.size = Pt(14)
            p_jp.font.name = "游明朝"
            p_jp.font.color.rgb = RGBColor(255, 255, 255)
            p_jp.alignment = PP_ALIGN.LEFT
            p_jp.space_before = Pt(4)

        # ────────────────────────────────────────────────────────
        # 🌟 ⑥ 内観ギャラリー（Gallery）の書き出し
        # ────────────────────────────────────────────────────────
        if "interior" in selected_pages:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            sw, sh = prs.slide_width, prs.slide_height
            luxury_gold = RGBColor(185, 160, 110)
            
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(20, 20, 25)
            
            gallery_bg = Image.new('RGB', (width, height), color=(20, 20, 25))
            preview_list.append({"type": "P.6 内観ギャラリー", "image": image_to_base64(gallery_bg)})

            if orientation == "landscape":
                en_y, jp_y, bar_y, main_y, main_h = Inches(0.2), Inches(0.7), Inches(1.4), Inches(1.6), Inches(2.5)
                sub1_y, sub2_y, sub_h = Inches(4.3), Inches(5.8), Inches(1.3)
                
                add_placeholder_box(slide, Inches(0.5), main_y, Inches(9.0), main_h, "【 メイン画像 】", luxury_gold, show_icon=True, label_text="MAIN SLOT")
                add_placeholder_box(slide, Inches(0.5), sub1_y, Inches(4.3), sub_h, "【 サブ画像1 】", luxury_gold, show_icon=True, label_text="SUB 1")
                add_placeholder_box(slide, Inches(5.2), sub1_y, Inches(4.3), sub_h, "【 サブ画像2 】", luxury_gold, show_icon=True, label_text="SUB 2")
                add_placeholder_box(slide, Inches(0.5), sub2_y, Inches(4.3), sub_h, "【 サブ画像3 】", luxury_gold, show_icon=True, label_text="SUB 3")
                add_placeholder_box(slide, Inches(5.2), sub2_y, Inches(4.3), sub_h, "【 サブ画像4 】", luxury_gold, show_icon=True, label_text="SUB 4")
            else:
                en_y, jp_y, bar_y, main_y, main_h = Inches(0.3), Inches(0.8), Inches(1.5), Inches(1.7), Inches(3.7)
                sub1_y, sub2_y, sub_h = Inches(5.6), Inches(7.6), Inches(1.8)
                
                add_placeholder_box(slide, Inches(0.5), main_y, Inches(6.5), main_h, "【 メイン画像 】", luxury_gold, show_icon=True, label_text="MAIN SLOT")
                add_placeholder_box(slide, Inches(0.5), sub1_y, Inches(3.1), sub_h, "【 サブ1 】", luxury_gold, show_icon=True, label_text="SUB 1")
                add_placeholder_box(slide, Inches(3.9), sub1_y, Inches(3.1), sub_h, "【 サブ2 】", luxury_gold, show_icon=True, label_text="SUB 2")
                add_placeholder_box(slide, Inches(0.5), sub2_y, Inches(3.1), sub_h, "【 サブ3 】", luxury_gold, show_icon=True, label_text="SUB 3")
                add_placeholder_box(slide, Inches(3.9), sub2_y, Inches(3.1), sub_h, "【 サブ4 】", luxury_gold, show_icon=True, label_text="SUB 4")

            gold_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, bar_y, sw, Pt(6))
            gold_bar.fill.solid()
            gold_bar.fill.fore_color.rgb = luxury_gold
            gold_bar.line.fill.background()

            tx_head_en = slide.shapes.add_textbox(sw * 0.05, en_y, sw * 0.5, Inches(0.5))
            tx_head_en.text_frame.paragraphs[0].text = "INTERIOR GALLERY"
            tx_head_en.text_frame.paragraphs[0].font.size, tx_head_en.text_frame.paragraphs[0].font.color.rgb = Pt(28), RGBColor(255, 255, 255)
            
            tx_head_jp = slide.shapes.add_textbox(sw * 0.05, jp_y, sw * 0.4, Inches(0.6))
            tx_head_jp.text_frame.paragraphs[0].text = "内観ギャラリー"
            tx_head_jp.text_frame.paragraphs[0].font.size, tx_head_jp.text_frame.paragraphs[0].font.name, tx_head_jp.text_frame.paragraphs[0].font.color.rgb = Pt(16), "游明朝", luxury_gold

        # ────────────────────────────────────────────────────────
        # 🌟 ⑦ 会社案内（Company）の書き出し
        # ────────────────────────────────────────────────────────
        if "company" in selected_pages:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            sw, sh = prs.slide_width, prs.slide_height
            luxury_gold = RGBColor(185, 160, 110)
            
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(20, 20, 25)
            
            company_bg = Image.new('RGB', (width, height), color=(20, 20, 25))
            preview_list.append({"type": "P.7 会社案内", "image": image_to_base64(company_bg)})

            if orientation == "landscape":
                main_y, main_h = Inches(0.4), Inches(4.2)
                sub_y, sub_h = Inches(4.8), Inches(2.1)
                sub_w = Inches(4.3)
                
                add_placeholder_box(slide, Inches(0.5), main_y, Inches(9.0), main_h, "【 会社案内 メイン画像 】", luxury_gold, show_icon=True, label_text="COMPANY MAIN")
                add_placeholder_box(slide, Inches(0.5), sub_y, sub_w, sub_h, "【 オフィス写真 】", luxury_gold, show_icon=True, label_text="OFFICE")
                add_placeholder_box(slide, Inches(5.2), sub_y, sub_w, sub_h, "【 スタッフ写真 】", luxury_gold, show_icon=True, label_text="STAFF")
            else:
                main_y, main_h = Inches(0.5), Inches(5.5)
                sub_y, sub_h = Inches(6.2), Inches(3.0)
                sub_w = Inches(3.1)
                
                add_placeholder_box(slide, Inches(0.5), main_y, Inches(6.5), main_h, "【 会社案内 メイン画像 】", luxury_gold, show_icon=True, label_text="COMPANY MAIN")
                add_placeholder_box(slide, Inches(0.5), sub_y, sub_w, sub_h, "【 オフィス写真 】", luxury_gold, show_icon=True, label_text="OFFICE")
                add_placeholder_box(slide, Inches(3.9), sub_y, sub_w, sub_h, "【 スタッフ写真 】", luxury_gold, show_icon=True, label_text="STAFF")

            tx_slogan = slide.shapes.add_textbox(0, sh * 0.94, sw, Inches(0.4))
            p_slogan = tx_slogan.text_frame.paragraphs[0]
            p_slogan.text = "「住まい」を通じて、お客様の人生に確かな価値を。"
            p_slogan.font.size, p_slogan.font.name = Pt(14), "游明朝"
            p_slogan.font.color.rgb = RGBColor(245, 240, 225)
            p_slogan.alignment = PP_ALIGN.CENTER

        # ────────────────────────────────────────────────────────
        # 🌟 ⑧ Contact & Thank you スライド（結び、共通で最後に追加）
        # ────────────────────────────────────────────────────────
        slide_thanks = prs.slides.add_slide(prs.slide_layouts[6]) 
        sw, sh = prs.slide_width, prs.slide_height
        gold_main = RGBColor(185, 160, 110)
        
        fill_thanks = slide_thanks.background.fill
        fill_thanks.solid()
        fill_thanks.fore_color.rgb = RGBColor(252, 248, 242)
        
        thanks_bg = Image.new('RGB', (width, height), color=(252, 248, 242))
        preview_list.append({"type": "P.8 結び", "image": image_to_base64(thanks_bg)})

        top_line = slide_thanks.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, sw, Pt(6))
        top_line.fill.solid()
        top_line.fill.fore_color.rgb = gold_main
        top_line.line.fill.background()

        tx_ty = slide_thanks.shapes.add_textbox(0, sh * 0.12, sw, Inches(1.0))
        p_ty = tx_ty.text_frame.paragraphs[0]
        p_ty.text = "THANK YOU"
        p_ty.font.size, p_ty.font.name, p_ty.alignment = Pt(48), "Arial", PP_ALIGN.CENTER

        map_box_w, map_box_h = Inches(5.5), Inches(2.2) 
        add_placeholder_box(slide_thanks, (sw - map_box_w) / 2, sh * 0.30, map_box_w, map_box_h, "店舗案内図（地図）")

        box_w, box_h, box_y = Inches(6.5), Inches(1.8), sh * 0.68 
        box = slide_thanks.shapes.add_shape(MSO_SHAPE.RECTANGLE, (sw - box_w) / 2, box_y, box_w, box_h)
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(255, 255, 255)
        box.line.color.rgb = RGBColor(240, 240, 240)
        box.line.width = Pt(1)

        tx_contact = slide_thanks.shapes.add_textbox((sw - box_w) / 2, box_y, box_w, box_h)
        tf_contact = tx_contact.text_frame
        tf_contact.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf_contact.margin_top = tf_contact.margin_bottom = 0
        tf_contact.clear()
        
        tel_info = extracted_info.get("tel", "0120-000-000")
        lic_info = extracted_info.get("license", "東京都知事（〇）第〇〇〇号")
        adr_info = extracted_info.get("address", "東京都...")

        p_tel = tf_contact.add_paragraph()
        p_tel.text = f"フリーダイヤル：{tel_info}"
        p_tel.font.size, p_tel.font.bold, p_tel.font.color.rgb, p_tel.alignment = Pt(26), True, gold_main, PP_ALIGN.CENTER

        p_info = tf_contact.add_paragraph()
        p_info.text = f"{lic_info}  | {adr_info}"
        p_info.font.size, p_info.font.color.rgb, p_info.alignment = Pt(10), RGBColor(120, 120, 120), PP_ALIGN.CENTER

        # 6. 完成したPowerPointファイルをバイナリ（データストリーム）としてBase64化
        pptx_io = BytesIO()
        prs.save(pptx_io)
        pptx_io.seek(0)
        pptx_base64 = base64.b64encode(pptx_io.getvalue()).decode('utf-8')

        print("✅ 【APB】すべてのAI解析・画像合成・パワポ作成がエラーなく完了しました！")
        return {
            "status": "success",
            "preview_images": preview_list,
            "pptx_base64": pptx_base64
        }
    except Exception as e:
        import traceback
        from fastapi.responses import JSONResponse
        
        error_details = traceback.format_exc()
        print(f"🔥 システムエラー発生:\n{error_details}")
        set_progress("❌ 処理中にエラーが発生しました")
        
        return JSONResponse(status_code=500, content={"detail": f"Pythonエラー詳細:\n{error_details}"})
@app.post("/generate_zumen_file")
def generate_zumen(
    title: str = Form(""),
    price: str = Form(""),
    address: str = Form(""),
    transport_line: str = Form(""), # 🌟 追加：路線名を受け取る
    transport_station: str = Form(""),
    transport_walk: str = Form(""),
    madori: str = Form(""),
    age: str = Form(""),
    right: str = Form(""),
    land_area: str = Form(""),
    building_area: str = Form(""),
    plan: str = Form(""),
    branch_name: str = Form("国分寺"),
    user_email: str = Form(""), # 🌟 ここを追加！ログインユーザーのメアドを受け取る
    design_num: str = Form("1"),
    full_summary: str = Form(""),
    footer_design_num: str = Form("1"), # 🌟 ここを追加！
    summary_font_size: str = Form("10"), # 🌟 ここを追加！WEBから送られてくる文字サイズを受け取る

    main_image: UploadFile = File(None),
    sub_image1: UploadFile = File(None),
    sub_image2: UploadFile = File(None),
    sub_image3: UploadFile = File(None),
    madori_image: UploadFile = File(None),
    tenpo_image: UploadFile = File(None),
    custom_footer_image: UploadFile = File(None), # 🌟 ここも追加！

    icon_image1: UploadFile = File(None),
    icon_image2: UploadFile = File(None),
    icon_image3: UploadFile = File(None),
    icon_image4: UploadFile = File(None),
    icon_image5: UploadFile = File(None),
    icon_image6: UploadFile = File(None)
):
    try:
        import os
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE # 🌟ここに追加！
        
        # 🌟【バグ修正】パワポで「_x000D_」になるエラーを完全に防ぐため、改行コードのゴミ(\r)を強制排除！
        full_summary = full_summary.replace('\r', '').replace('_x000D_', '')
        title = title.replace('\r', '').replace('_x000D_', '')
        address = address.replace('\r', '').replace('_x000D_', '')
        transport_line = transport_line.replace('\r', '').replace('_x000D_', '') # 🌟 追加
        
        print(f"🚀 【販売図面】デザイン{design_num}のパワポ生成を開始します...")

        # 1. パワポの土台を作成（横長スライド）
        prs = Presentation()
        prs.slide_width, prs.slide_height = Inches(10), Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 2. 店舗データの取得（🌟 ログインユーザーに紐づく店舗をDBから自動判定！）
        db = SessionLocal()
        try:
            # リクエストしてきたユーザーの情報を取得
            db_user = db.query(User).filter(User.email == user_email).first()
            
            if db_user:
                target_branch = db_user.branch_key
            else:
                target_branch = branch_name

            # 判定された店舗キーで店舗情報を取得
            db_branch = db.query(Branch).filter(Branch.branch_key == target_branch).first()
            if not db_branch:
                db_branch = db.query(Branch).filter(Branch.branch_key == "国分寺").first()
            
            branch = {
                "full_name": db_branch.full_name,
                "license": db_branch.license,
                "address": db_branch.address,
                "tel": db_branch.tel
            }
        finally:
            db.close()

        # 共通の描画ヘルパー関数
        def add_color_box(left, top, width, height, text, bg_color, font_size=14):
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = bg_color
            shape.line.fill.background()
            tf = shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(font_size)
            p.font.name = "游明朝"
            p.font.color.rgb = RGBColor(80, 80, 80)
            p.alignment = PP_ALIGN.CENTER
           # 🌟🌟【新規追加】画像がアップロードされていれば写真を、なければ色枠を配置する賢い関数
        def add_smart_image(left, top, width, height, upload_file, fallback_text, bg_color, font_size=14):
            if upload_file and upload_file.filename:
                import tempfile
                import os
                try:
                    suffix = os.path.splitext(upload_file.filename)[1]
                    if not suffix: suffix = ".png"
                    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
                        tmp.write(upload_file.file.read())
                        tmp_path = tmp.name
                    
                    # 枠の大きさに合わせて画像を自動リサイズして配置
                    slide.shapes.add_picture(tmp_path, left, top, width=width, height=height)
                    os.remove(tmp_path)
                    
                    # 🌟 追加：画像の右下に「画像1」などのラベル（キャプション）を配置する魔法
                    lbl_w, lbl_h = Inches(0.8), Inches(0.25)
                    lbl_l = left + width - lbl_w
                    lbl_t = top + height - lbl_h
                    
                    lbl_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, lbl_l, lbl_t, lbl_w, lbl_h)
                    lbl_bg.fill.solid()
                    lbl_bg.fill.fore_color.rgb = RGBColor(128, 154, 185) # 理想の画像に合わせたブルーグレー
                    try: lbl_bg.fill.transparency = 0.1 # 10%だけ透けさせて背景の画像になじませる
                    except: pass
                    lbl_bg.line.fill.background() # 枠線は消す
                    
                    tf = lbl_bg.text_frame
                    tf.margin_top = tf.margin_bottom = tf.margin_left = tf.margin_right = 0
                    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                    p = tf.paragraphs[0]
                    p.text = fallback_text # 「画像1」などのテキストをそのまま入れる
                    p.font.size = Pt(10)
                    p.font.color.rgb = RGBColor(255, 255, 255) # 白色
                    p.font.name = "游ゴシック"
                    p.alignment = PP_ALIGN.CENTER

                    return
                except Exception as e:
                    print(f"画像配置エラー: {e}")
            
            # 画像がない場合やエラー時は今まで通りの色枠を出す
            add_color_box(left, top, width, height, fallback_text, bg_color, font_size) 
# 🌟 追加：グラデーション背景を自動生成して滑らかにフェードさせる魔法の関数（フェード開始・終了位置の完全指定版！）
        def add_gradient_box(left, top, w, h, color1, color2, start_pos=0.0, end_pos=1.0, horizontal=False):
            import io
            from PIL import Image, ImageDraw
            img = Image.new('RGB', (300, 2) if horizontal else (2, 300))
            draw = ImageDraw.Draw(img)
            start_pixel = int(300 * start_pos)
            end_pixel = int(300 * end_pos)
            
            for i in range(300):
                if i <= start_pixel:
                    # 指定した開始位置までは、color1（青）をベタ塗り
                    r, g, b = color1
                elif i >= end_pixel:
                    # 指定した終了位置からは、color2（白）をベタ塗り
                    r, g, b = color2
                else:
                    # その間だけを滑らかにフェードアウト！
                    ratio = (i - start_pixel) / (end_pixel - start_pixel)
                    r = int(color1[0] + (color2[0] - color1[0]) * ratio)
                    g = int(color1[1] + (color2[1] - color1[1]) * ratio)
                    b = int(color1[2] + (color2[2] - color1[2]) * ratio)
                
                if horizontal:
                    draw.line([(i, 0), (i, 1)], fill=(r,g,b))
                else:
                    draw.line([(0, i), (1, i)], fill=(r,g,b))
            img_io = io.BytesIO()
            img.save(img_io, format='PNG')
            img_io.seek(0)
            slide.shapes.add_picture(img_io, left, top, width=w, height=h)

        # ==========================================
        # 🎨 レイアウト分岐処理
        # ==========================================
        box_color = RGBColor(253, 232, 215) # 肌色
        pink_color = RGBColor(248, 232, 248) # ピンク
        gray_line = RGBColor(120, 120, 120)
        footer_y = Inches(6.65) # フッターの基本Y座標

        if design_num == "2":
            # ----------------------------------------------------
            # 🟦 デザイン2（白×水色 スッキリレイアウト）
            # ----------------------------------------------------
            # 全体の白枠
            bg_frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.1), Inches(0.1), Inches(9.8), Inches(7.3))
            bg_frame.fill.solid()
            bg_frame.fill.fore_color.rgb = RGBColor(255, 255, 255)
            bg_frame.line.color.rgb = RGBColor(207, 216, 220)
            
            # 左側：画像1
            add_smart_image(Inches(0.3), Inches(0.3), Inches(4.2), Inches(3.8), main_image, "画像1", box_color, 24)
            
            # タイトル
            tb_title = slide.shapes.add_textbox(Inches(0.3), Inches(4.2), Inches(2.6), Inches(0.8))
            tf_title = tb_title.text_frame
            tf_title.vertical_anchor = MSO_ANCHOR.MIDDLE  # 🌟 追加：文字を上下の中央に美しく揃える！
            tf_title.word_wrap = True  # 🌟 追加：枠の右端で自動的に折り返す
            tf_title.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE  # 🌟 追加：文字が多すぎる場合は自動縮小
            p_title = tf_title.paragraphs[0]
            p_title.text = title if title else "タイトルを入力"
            p_title.font.size = Pt(28)
            p_title.font.name = "游明朝"
            p_title.font.bold = True
            
            # 価格（Runを使って文字ごとにサイズを変える魔法）
            tb_price = slide.shapes.add_textbox(Inches(2.5), Inches(4.25), Inches(2.0), Inches(0.8))
            tf_price = tb_price.text_frame
            tf_price.margin_bottom = tf_price.margin_top = 0
            p_price = tf_price.paragraphs[0]
            p_price.alignment = PP_ALIGN.RIGHT
            
            r_label = p_price.add_run()
            r_label.text = "販売価格 "
            r_label.font.size = Pt(10)
            r_label.font.name = "游明朝"
            r_label.font.color.rgb = RGBColor(100, 100, 100)
            
            r_val = p_price.add_run()
            r_val.text = price
            r_val.font.size = Pt(32) # 🌟金額だけ特大に
            r_val.font.name = "游明朝"
            r_val.font.bold = True
            
            r_unit = p_price.add_run()
            r_unit.text = " 万円"
            r_unit.font.size = Pt(14)
            r_unit.font.name = "游明朝"

            # 設備アイコン
            tb_eq = slide.shapes.add_textbox(Inches(0.3), Inches(5.1), Inches(0.6), Inches(0.5))
            tb_eq.text_frame.paragraphs[0].text = "物件\n設備"
            tb_eq.text_frame.paragraphs[0].font.size = Pt(10)
            tb_eq.text_frame.paragraphs[0].font.name = "游明朝"
            
            # 🌟 変更：アイコン画像があれば貼り付ける処理（デザイン2）
            icons = [icon_image1, icon_image2, icon_image3, icon_image4, icon_image5, icon_image6]
            import tempfile
            from PIL import Image

            for i in range(6):
                current_left = Inches(1.0 + i*0.6)
                current_top = Inches(5.1)
                icon_file = icons[i]
                
                if icon_file and icon_file.filename:
                    suffix = os.path.splitext(icon_file.filename)[1]
                    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
                        tmp.write(icon_file.file.read())
                        tmp_path = tmp.name
                    
                    with Image.open(tmp_path) as img:
                        img_w, img_h = img.size
                    
                    # 高さを枠の高さ(0.5インチ)に合わせる
                    target_h_inches = 0.5
                    target_w_inches = target_h_inches * (img_w / img_h)
                    
                    # 枠の横幅(0.5インチ)との差分から、中央寄せするためのX座標を計算
                    offset_x_inches = (0.5 - target_w_inches) / 2
                    centered_left = current_left + Inches(offset_x_inches)
                    
                    slide.shapes.add_picture(tmp_path, centered_left, current_top, height=Inches(target_h_inches))
                    os.remove(tmp_path)
                else:
                    add_color_box(current_left, current_top, Inches(0.5), Inches(0.5), "設備", pink_color, 9)

            # 🌟 変更（右側：画像2,3,4）
            add_smart_image(Inches(4.8), Inches(0.3), Inches(1.5), Inches(1.4), sub_image1, "画像2", box_color, 14)
            add_smart_image(Inches(6.4), Inches(0.3), Inches(1.5), Inches(1.4), sub_image2, "画像3", box_color, 14)
            add_smart_image(Inches(8.0), Inches(0.3), Inches(1.5), Inches(1.4), sub_image3, "画像4", box_color, 14)

            # 交通アクセス (上の水色線)
            blue_line1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.8), Inches(1.9), Inches(4.7), Pt(1))
            blue_line1.fill.solid()
            blue_line1.fill.fore_color.rgb = RGBColor(127, 179, 213)
            blue_line1.line.fill.background()
            
            # 🌟 「A c c e s s」の文字（白背景にして線の上に重ねる）
            tb_acc_label = slide.shapes.add_textbox(Inches(4.9), Inches(1.78), Inches(1.0), Inches(0.25))
            tb_acc_label.fill.solid()
            tb_acc_label.fill.fore_color.rgb = RGBColor(255, 255, 255) # 背景を白くして線を隠す
            tb_acc_label.line.fill.background()
            p_acc_label = tb_acc_label.text_frame.paragraphs[0]
            p_acc_label.text = "A c c e s s"
            p_acc_label.font.size = Pt(9)
            p_acc_label.font.name = "Arial"
            p_acc_label.font.color.rgb = RGBColor(84, 153, 199)

            # 交通テキスト（こちらもサイズに強弱をつける）
            tb_acc = slide.shapes.add_textbox(Inches(4.8), Inches(2.0), Inches(4.7), Inches(0.5))
            p_acc = tb_acc.text_frame.paragraphs[0]
            p_acc.alignment = PP_ALIGN.CENTER
            
            r_acc1 = p_acc.add_run()
            r_acc1.text = "交通 "
            r_acc1.font.size = Pt(13)
            r_acc1.font.name = "游明朝"
            
            # 🌟 カッコが二重になるのを防ぐ（「」を消してから付け直す）
            clean_station = transport_station.replace("「", "").replace("」", "")
            r_acc2 = p_acc.add_run()
            r_acc2.text = f"「{clean_station}」"
            r_acc2.font.size = Pt(20)
            r_acc2.font.name = "游明朝"
            r_acc2.font.bold = True
            
            r_acc3 = p_acc.add_run()
            r_acc3.text = f" 徒歩 {transport_walk} 分"
            r_acc3.font.size = Pt(13)
            r_acc3.font.name = "游明朝"
            
            # 交通アクセス (下の水色線)
            blue_line2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.8), Inches(2.6), Inches(4.7), Pt(1))
            blue_line2.fill.solid()
            blue_line2.fill.fore_color.rgb = RGBColor(127, 179, 213)
            blue_line2.line.fill.background()

            # 間取り図
            add_smart_image(Inches(4.8), Inches(2.8), Inches(4.7), Inches(2.8), madori_image, "間取り図", box_color, 20)

            # 水色帯（物件詳細）※高さを0.95インチに固定し、枠が下に伸びないようにする
            blue_band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.11), Inches(5.7), Inches(9.78), Inches(0.95))
            blue_band.fill.solid()
            blue_band.fill.fore_color.rgb = RGBColor(221, 240, 248) # #DDF0F8
            blue_band.line.fill.background()

            # 「■ 物件詳細情報」の見出しをパワポにも描画する
            tb_summary_title = slide.shapes.add_textbox(Inches(0.2), Inches(5.75), Inches(9.6), Inches(0.2))
            tf_summary_title = tb_summary_title.text_frame
            tf_summary_title.margin_top = tf_summary_title.margin_bottom = tf_summary_title.margin_left = tf_summary_title.margin_right = 0
            p_summary_title = tf_summary_title.paragraphs[0]
            p_summary_title.text = "■ 物件詳細情報"
            p_summary_title.font.size = Pt(10)
            p_summary_title.font.bold = True
            p_summary_title.font.color.rgb = RGBColor(84, 153, 199) # 画面と同じ青色(#5499C7)
            p_summary_title.font.name = "游ゴシック"

            # 🌟 受け取った全項目を枠内に収めるため、6列にしてフォントサイズと行間を調整する
            items = [item.strip() for item in full_summary.split('|||') if item.strip()]
            cols = 6  # 🌟 5列から6列に変更（1列あたりの行数を減らしてあふれを防ぐ）
            box_width = Inches(9.6) / cols
            items_per_col = (len(items) + cols - 1) // cols
            
            for i in range(cols):
                col_items = items[i * items_per_col : (i + 1) * items_per_col]
                if not col_items: continue
                
                # 見出しの分だけ、項目のスタート位置（Y座標）を少し下げる
                tb_info = slide.shapes.add_textbox(Inches(0.2) + i * box_width, Inches(5.95), box_width, Inches(0.65))
                tf_info = tb_info.text_frame
                
                # 🌟【案2を採用】枠の幅は固定し、改行を許可（True）した上で自動縮小する！
                tf_info.word_wrap = True 
                tf_info.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE 
                
                # 右側に少し余白を設けて、隣の列の文字と絶対にくっつかないようにする
                tf_info.margin_top = tf_info.margin_bottom = tf_info.margin_left = 0
                tf_info.margin_right = Inches(0.05)
                
                for idx, item_text in enumerate(col_items):
                    p = tf_info.paragraphs[0] if idx == 0 else tf_info.add_paragraph()
                    p.text = item_text
                    
                    p.font.size = Pt(7.0) 
                    p.font.name = "游明朝"
                    p.font.color.rgb = RGBColor(80, 80, 80) 
                    p.line_spacing = Pt(8.5)
                    
                    # 🌟【究極の整列魔法】パワポの裏側（XML）を直接操作してインデントを設定！
                    indent_emu = int(Pt(7.0 * 8.5))
                    pPr = p._p.get_or_add_pPr()
                    pPr.set('marL', str(indent_emu))
                    pPr.set('indent', str(-indent_emu))
            
            # 🌟 フッター位置はスライド内に収まる基本位置（6.65インチ）に戻す
            footer_y = Inches(6.65)

            # フッターの上の濃い線を引く
            f_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.1), footer_y, Inches(9.8), Pt(2))
            f_line.fill.solid()
            f_line.fill.fore_color.rgb = RGBColor(30, 45, 61)
            f_line.line.fill.background()

        elif design_num == "3":
            # --- ▼ 左側エリア ▼ ---
            cyan_color = RGBColor(86, 180, 203)
            # 🌟 変更（上段：画像2, 3, 4）
            add_smart_image(Inches(0.2), Inches(0.2), Inches(2.1), Inches(1.3), sub_image1, "画像2", box_color, 14)
            add_smart_image(Inches(2.4), Inches(0.2), Inches(2.1), Inches(1.3), sub_image2, "画像3", box_color, 14)
            add_smart_image(Inches(4.6), Inches(0.2), Inches(2.1), Inches(1.3), sub_image3, "画像4", box_color, 14)
            
            # 中段：タイトル＆交通アクセスの水色帯
            cyan_band_left = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.2), Inches(1.6), Inches(6.5), Inches(1.0))
            cyan_band_left.fill.solid()
            cyan_band_left.fill.fore_color.rgb = cyan_color
            cyan_band_left.line.fill.background()
            
            # タイトル
            tb_title = slide.shapes.add_textbox(Inches(0.2), Inches(1.6), Inches(3.7), Inches(1.0))
            tf_title = tb_title.text_frame
            tf_title.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf_title.word_wrap = True  # 🌟 追加：枠の右端で自動的に折り返す
            tf_title.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE  # 🌟 追加：文字が多すぎる場合は自動縮小
            p_title = tf_title.paragraphs[0]
            p_title.text = title if title else "タイトルを入力"
            p_title.font.size = Pt(22)
            p_title.font.bold = True
            p_title.font.color.rgb = RGBColor(255, 255, 255)
            p_title.font.name = "游明朝"
            p_title.alignment = PP_ALIGN.CENTER
            
            # 縦の白線（タイトルと交通の区切り）
            v_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.9), Inches(1.7), Pt(1), Inches(0.8))
            v_line.fill.solid()
            v_line.fill.fore_color.rgb = RGBColor(255, 255, 255)
            v_line.line.fill.background()
            
            # ACCESS 縦文字
            tb_acc_v = slide.shapes.add_textbox(Inches(3.95), Inches(1.65), Inches(0.2), Inches(0.9))
            tf_acc_v = tb_acc_v.text_frame
            # 🌟 0 ではなく Inches(0) と指定する
            tf_acc_v.margin_left = tf_acc_v.margin_right = Inches(0) 
            p_acc_v = tf_acc_v.paragraphs[0]
            p_acc_v.text = "A\nC\nC\nE\nS\nS"
            p_acc_v.font.size = Pt(7)
            p_acc_v.font.color.rgb = RGBColor(255, 255, 255)
            p_acc_v.alignment = PP_ALIGN.CENTER
            # 🌟 0.9 ではなく Pt(8) などの正規の単位にする
            p_acc_v.line_spacing = Pt(8)
            
            # 交通情報
            tb_acc = slide.shapes.add_textbox(Inches(4.15), Inches(1.6), Inches(2.55), Inches(1.0))
            tf_acc = tb_acc.text_frame
            tf_acc.margin_left = tf_acc.margin_right = 0
            tf_acc.vertical_anchor = MSO_ANCHOR.MIDDLE
            p_acc = tf_acc.paragraphs[0]
            p_acc.alignment = PP_ALIGN.CENTER
            
            r_acc1 = p_acc.add_run()
            r_acc1.text = "交通  "
            r_acc1.font.size = Pt(10)
            r_acc1.font.color.rgb = RGBColor(255, 255, 255)
            r_acc1.font.name = "游明朝"
            
            clean_station = transport_station.replace("「", "").replace("」", "")
            r_acc2 = p_acc.add_run()
            r_acc2.text = f"「{clean_station}」"
            r_acc2.font.size = Pt(14)
            r_acc2.font.bold = True
            r_acc2.font.color.rgb = RGBColor(255, 255, 255)
            r_acc2.font.name = "游明朝"
            
            r_acc3 = p_acc.add_run()
            r_acc3.text = f"  徒歩 {transport_walk} 分"
            r_acc3.font.size = Pt(10)
            r_acc3.font.color.rgb = RGBColor(255, 255, 255)
            r_acc3.font.name = "游明朝"

            # 🌟 変更（下段：メイン画像＆間取り図）
            add_smart_image(Inches(0.2), Inches(2.7), Inches(3.2), Inches(3.2), main_image, "画像1", box_color, 24)
            add_smart_image(Inches(3.5), Inches(2.7), Inches(3.2), Inches(3.2), madori_image, "間取り図", box_color, 24)
            
            # 設備アイコン
            tb_eq = slide.shapes.add_textbox(Inches(0.2), Inches(6.0), Inches(0.8), Inches(0.5))
            p_eq = tb_eq.text_frame.paragraphs[0]
            p_eq.text = "物件\n設備"
            p_eq.font.size = Pt(10)
            p_eq.font.color.rgb = RGBColor(80, 80, 80)
            p_eq.font.name = "游明朝"
            
            icons = [icon_image1, icon_image2, icon_image3, icon_image4, icon_image5, icon_image6]
            import tempfile
            from PIL import Image

            for i in range(6):
                current_left = Inches(1.0 + i*0.95)
                current_top = Inches(6.0)
                icon_file = icons[i]
                
                if icon_file and icon_file.filename:
                    suffix = os.path.splitext(icon_file.filename)[1]
                    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
                        tmp.write(icon_file.file.read())
                        tmp_path = tmp.name
                    
                    with Image.open(tmp_path) as img:
                        img_w, img_h = img.size
                    
                    target_h_inches = 0.5
                    target_w_inches = target_h_inches * (img_w / img_h)
                    offset_x_inches = (0.85 - target_w_inches) / 2
                    centered_left = current_left + Inches(offset_x_inches)
                    
                    slide.shapes.add_picture(tmp_path, centered_left, current_top, height=Inches(target_h_inches))
                    os.remove(tmp_path)
                else:
                    add_color_box(current_left, current_top, Inches(0.85), Inches(0.5), "設備", pink_color, 10)

            # --- ▼ 右側エリア（水色背景の物件詳細・価格） ▼ ---
            cyan_band_right = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(0.2), Inches(3.0), Inches(6.3))
            cyan_band_right.fill.solid()
            cyan_band_right.fill.fore_color.rgb = cyan_color
            cyan_band_right.line.fill.background()
            
            # 価格エリア
            tb_price_label = slide.shapes.add_textbox(Inches(6.9), Inches(0.6), Inches(0.6), Inches(0.6))
            p_pr_l = tb_price_label.text_frame.paragraphs[0]
            p_pr_l.text = "販売\n価格"
            p_pr_l.font.size = Pt(10)
            p_pr_l.font.color.rgb = RGBColor(255, 255, 255)
            p_pr_l.font.name = "游明朝"
            
            tb_price = slide.shapes.add_textbox(Inches(7.5), Inches(0.4), Inches(2.2), Inches(0.8))
            p_price = tb_price.text_frame.paragraphs[0]
            p_price.alignment = PP_ALIGN.RIGHT
            
            r_val = p_price.add_run()
            r_val.text = price
            r_val.font.size = Pt(36)
            r_val.font.bold = True
            r_val.font.color.rgb = RGBColor(255, 255, 255)
            r_val.font.name = "游明朝"
            
            r_unit = p_price.add_run()
            r_unit.text = " 万円"
            r_unit.font.size = Pt(14)
            r_unit.font.color.rgb = RGBColor(255, 255, 255)
            r_unit.font.name = "游明朝"

            # 価格下の白線
            white_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.0), Inches(1.3), Inches(2.6), Pt(1))
            white_line.fill.solid()
            white_line.fill.fore_color.rgb = RGBColor(255, 255, 255)
            white_line.line.fill.background()
            
            # 🌟【ここを追加】「■ 物件詳細情報」の見出しをパワポにも描画する
            tb_summary_title = slide.shapes.add_textbox(Inches(6.9), Inches(1.45), Inches(2.8), Inches(0.2))
            tf_summary_title = tb_summary_title.text_frame
            tf_summary_title.margin_top = tf_summary_title.margin_bottom = tf_summary_title.margin_left = tf_summary_title.margin_right = 0
            p_summary_title = tf_summary_title.paragraphs[0]
            p_summary_title.text = "■ 物件詳細情報"
            p_summary_title.font.size = Pt(11)
            p_summary_title.font.bold = True
            p_summary_title.font.color.rgb = RGBColor(255, 255, 255)
            p_summary_title.font.name = "游ゴシック"

            # 🌟 変更（物件詳細情報を2列・自動縮小ではみ出し防止！）
            items = [item.strip() for item in full_summary.split('|||') if item.strip()]
            cols = 2
            box_width = Inches(2.8) / cols
            items_per_col = (len(items) + cols - 1) // cols
            
            for i in range(cols):
                col_items = items[i * items_per_col : (i + 1) * items_per_col]
                if not col_items: continue
                
                # 🌟 見出しを追加した分、リストの開始位置（Y座標）を少し下げる（Inches(1.7)に変更）
                tb_info = slide.shapes.add_textbox(Inches(6.9) + i * box_width, Inches(1.7), box_width, Inches(4.6))
                tf_info = tb_info.text_frame
                tf_info.word_wrap = True 
                tf_info.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE 
                tf_info.margin_top = tf_info.margin_bottom = tf_info.margin_left = Inches(0)
                tf_info.margin_right = Inches(0.05)
                
                for idx, item_text in enumerate(col_items):
                    p = tf_info.paragraphs[0] if idx == 0 else tf_info.add_paragraph()
                    p.text = item_text
                    
                    p.font.size = Pt(8.5)
                    p.font.name = "游明朝"
                    p.font.color.rgb = RGBColor(255, 255, 255) 
                    p.line_spacing = Pt(12.0)
                    p.space_before = Pt(6.0)
                    
                    # 🌟【究極の整列魔法】パワポの裏側（XML）を直接操作してインデントを設定！
                    indent_emu = int(Pt(8.5 * 8.5))
                    pPr = p._p.get_or_add_pPr()
                    pPr.set('marL', str(indent_emu))
                    pPr.set('indent', str(-indent_emu))
                
            # フッターの線の開始位置
            footer_y = Inches(6.65)
                
            # フッターの線の開始位置
            footer_y = Inches(6.65)
            # フッターの上の線を引く
            f_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.2), footer_y - Inches(0.05), Inches(9.6), Pt(1.5))
            f_line.fill.solid()
            f_line.fill.fore_color.rgb = RGBColor(120, 120, 120)
            f_line.line.fill.background()

        else:
            # ----------------------------------------------------
            # 🟦 デザイン1（最新：モダンな背景色分けレイアウト）
            # ----------------------------------------------------
            
            # 🌟 1. 上部の背景を、以前のしっかりした青色に戻し、28%〜40%で白へフェードアウト！
            add_gradient_box(0, 0, Inches(10), Inches(1.4), (170, 200, 230), (250, 252, 255), start_pos=0.28, end_pos=0.40, horizontal=True)

            # 🌟 2. 下の画像エリアも、上部と「全く同じ横グラデーション（幅も10インチ）」に広げることで、段差を完全に消しつつ画像間で美しくフェードさせます！
            add_gradient_box(0, Inches(1.4), Inches(10), Inches(4.5), (170, 200, 230), (250, 252, 255), start_pos=0.28, end_pos=0.40, horizontal=True)

            # 🌟 3. 下の設備エリアを「左(水色)から右(白)へのグラデーション」に
            add_gradient_box(0, Inches(5.9), Inches(10), Inches(0.75), (230, 240, 250), (255, 255, 255), start_pos=0.0, end_pos=1.0, horizontal=True)

            # 4. 右側の物件詳細用の白い角丸カード
            bg_right = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.25), Inches(1.45), Inches(2.65), Inches(4.4))
            bg_right.fill.solid()
            bg_right.fill.fore_color.rgb = RGBColor(255, 255, 255)
            bg_right.line.fill.background()
            bg_right.adjustments[0] = 0.05 
            
            # 🌟 5. 白いカードに「ふんわりとした影（ドロップシャドウ）」をつけて高級感をアップ！
            try:
                bg_right.shadow.inherit = False
                bg_right.shadow.color.rgb = RGBColor(0, 0, 0)
                bg_right.shadow.alpha = 0.85     # 影の薄さ（15%の濃さ）
                bg_right.shadow.blur_radius = Pt(10) # 影のぼかし具合
                bg_right.shadow.distance = Pt(2)     # 影の落ちる距離
                bg_right.shadow.angle = 45           # 影の角度（右下）
            except:
                pass 

            base_dir = os.path.dirname(__file__)
            tenpo_img_path = os.path.join(base_dir, "static", "tenpo.png")
            
            if os.path.exists(tenpo_img_path):
                slide.shapes.add_picture(tenpo_img_path, Inches(0.2), Inches(0.2), width=Inches(2.8), height=Inches(1.1))
            else:
                add_color_box(Inches(0.2), Inches(0.2), Inches(2.8), Inches(1.1), "店舗写真", RGBColor(235, 235, 235), 12)
            
            # 店舗写真の上に重なるテキスト
            tb_tenpo = slide.shapes.add_textbox(Inches(0.2), Inches(0.2), Inches(2.8), Inches(1.1))
            tf_tenpo = tb_tenpo.text_frame
            tf_tenpo.vertical_anchor = MSO_ANCHOR.MIDDLE
            p_tenpo = tf_tenpo.paragraphs[0]
            p_tenpo.text = "「住まい」のもっと先へ。"
            p_tenpo.font.size = Pt(14)
            p_tenpo.font.bold = True
            p_tenpo.font.color.rgb = RGBColor(255, 255, 255)
            p_tenpo.font.name = "游ゴシック"
            p_tenpo.alignment = PP_ALIGN.CENTER
            try: p_tenpo.font.shadow = True
            except: pass

            # 🌟 追加：タイトルの上下に高級感のある飾り線を復活させる魔法
            line_color = RGBColor(170, 200, 230) # 上品な薄いブルーグレー
            line_top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.2), Inches(0.25), Inches(4.0), Pt(1.5))
            line_top.fill.solid()
            line_top.fill.fore_color.rgb = line_color
            line_top.line.fill.background()
            
            line_bottom = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.2), Inches(1.25), Inches(4.0), Pt(1.5))
            line_bottom.fill.solid()
            line_bottom.fill.fore_color.rgb = line_color
            line_bottom.line.fill.background()

            # タイトル
            tb_title = slide.shapes.add_textbox(Inches(3.2), Inches(0.35), Inches(4.0), Inches(0.9))
            tf_title = tb_title.text_frame
            tf_title.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf_title.word_wrap = True
            tf_title.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            p_title = tf_title.paragraphs[0]
            p_title.text = title if title else "タイトルを入力"
            p_title.font.size = Pt(26)
            p_title.font.bold = True
            p_title.font.color.rgb = RGBColor(16, 51, 93) 
            p_title.alignment = PP_ALIGN.CENTER
            
            # 「販売/価格」のラベル
            tb_price_lbl = slide.shapes.add_textbox(Inches(6.8), Inches(0.55), Inches(0.6), Inches(0.6))
            tf_lbl = tb_price_lbl.text_frame
            tf_lbl.margin_top = tf_lbl.margin_bottom = tf_lbl.margin_left = tf_lbl.margin_right = 0
            p_lbl = tf_lbl.paragraphs[0]
            p_lbl.text = "販売\n価格"
            p_lbl.font.size = Pt(11)
            p_lbl.font.name = "游明朝"
            p_lbl.font.color.rgb = RGBColor(100, 100, 100)
            p_lbl.alignment = PP_ALIGN.RIGHT
            p_lbl.line_spacing = Pt(13)

            # 金額部分
            tb_price_val = slide.shapes.add_textbox(Inches(7.4), Inches(0.35), Inches(2.4), Inches(0.8))
            tf_val = tb_price_val.text_frame
            tf_val.margin_top = tf_val.margin_bottom = tf_val.margin_left = tf_val.margin_right = 0
            tf_val.vertical_anchor = MSO_ANCHOR.MIDDLE
            p_val = tf_val.paragraphs[0]
            p_val.alignment = PP_ALIGN.RIGHT
            
            r_val = p_val.add_run()
            r_val.text = price
            r_val.font.size = Pt(40)
            r_val.font.name = "游明朝"
            r_val.font.bold = True
            r_val.font.color.rgb = RGBColor(16, 51, 93) 
            
            r_unit = p_val.add_run()
            r_unit.text = " 万円"
            r_unit.font.size = Pt(16)
            r_unit.font.name = "游明朝"
            r_unit.font.color.rgb = RGBColor(16, 51, 93)

            # 🌟 5. 画像枠の代替色（画像がない時の空箱の色）を肌色から「真っ白（255,255,255）」に変更
            add_smart_image(Inches(0.2), Inches(1.5), Inches(3.4), Inches(2.8), main_image, "画像1", RGBColor(255, 255, 255), 18)
            add_smart_image(Inches(3.7), Inches(1.5), Inches(3.4), Inches(2.8), madori_image, "間取り図", RGBColor(255, 255, 255), 18)
            add_smart_image(Inches(0.2), Inches(4.4), Inches(2.2), Inches(1.4), sub_image1, "画像2", RGBColor(255, 255, 255), 16)
            add_smart_image(Inches(2.5), Inches(4.4), Inches(2.2), Inches(1.4), sub_image2, "画像3", RGBColor(255, 255, 255), 16)
            add_smart_image(Inches(4.8), Inches(4.4), Inches(2.2), Inches(1.4), sub_image3, "画像4", RGBColor(255, 255, 255), 16)

            line_v = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.2), Inches(1.4), Pt(1), Inches(4.5))
            line_v.fill.solid()
            line_v.fill.fore_color.rgb = gray_line
            line_v.line.fill.background()

            # 🌟 修正：Web画面と同じように2段構えにし、駅名を大きく表示する魔法
            tb_acc = slide.shapes.add_textbox(Inches(7.3), Inches(1.45), Inches(2.5), Inches(0.65))
            tf_acc = tb_acc.text_frame
            tf_acc.clear()
            tf_acc.margin_top = tf_acc.margin_bottom = 0
            tf_acc.vertical_anchor = MSO_ANCHOR.MIDDLE  # 🌟 追加：文字を上下の中央に美しく揃える！
            
            # 1行目：路線名
            p_line = tf_acc.paragraphs[0]
            p_line.text = transport_line if transport_line else "交通"
            p_line.font.size = Pt(10)
            p_line.font.color.rgb = RGBColor(100, 100, 100)
            p_line.font.name = "游明朝"
            p_line.alignment = PP_ALIGN.CENTER
            
            # 2行目：駅名と徒歩分数
            p_sta = tf_acc.add_paragraph()
            clean_station = transport_station.replace("「", "").replace("」", "")
            
            r_sta = p_sta.add_run()
            r_sta.text = f"{clean_station} "
            r_sta.font.size = Pt(18)
            r_sta.font.bold = True
            r_sta.font.color.rgb = RGBColor(80, 80, 80)
            r_sta.font.name = "游明朝"
            
            r_walk = p_sta.add_run()
            r_walk.text = f"徒歩 {transport_walk} 分"
            r_walk.font.size = Pt(10)
            r_walk.font.color.rgb = RGBColor(100, 100, 100)
            r_walk.font.name = "游明朝"
            p_sta.alignment = PP_ALIGN.CENTER
            
            line_acc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.3), Inches(2.1), Inches(2.5), Pt(1))
            line_acc.fill.solid()
            line_acc.fill.fore_color.rgb = gray_line
            line_acc.line.fill.background()

            tb_summ_t = slide.shapes.add_textbox(Inches(7.3), Inches(2.2), Inches(2.5), Inches(0.4))
            p_summ_t = tb_summ_t.text_frame.paragraphs[0]
            p_summ_t.text = "■ 物件詳細情報"
            p_summ_t.font.size = Pt(10)
            p_summ_t.font.bold = True

            items = [item.strip() for item in full_summary.split('|||') if item.strip()]
            
            tb_info = slide.shapes.add_textbox(Inches(7.3), Inches(2.5), Inches(2.5), Inches(3.45))
            tf_info = tb_info.text_frame
            tf_info.word_wrap = True 
            # 🚨 パワポの勝手な縮小機能をオフにして、こちらの計算を絶対優先させます
            tf_info.margin_top = tf_info.margin_bottom = tf_info.margin_left = tf_info.margin_right = 0
            
            # 🌟 【等間隔＆ジャストフィットの究極魔法】
            available_height_pt = 240.0 # 枠の高さ(3.45インチ)の実質利用可能ポイント
            
            # まず、折り返しを考慮した「実質的な行数」を見積もる
            num_lines = 0
            for item_text in items:
                # 🌟 修正：ユーザーが入力した改行（\n）も1行として正確にカウントする
                for line in item_text.split('\n'):
                    if len(line) > 19: # 20文字以上は2行に折り返されると判定
                        num_lines += 2
                    else:
                        num_lines += 1

            if num_lines == 0: num_lines = 1
            
            # ギリギリまで文字を大きく（最大9.5Pt、最小5.0Pt）
            calc_font_size = (available_height_pt / num_lines) * 0.8
            ppt_font_size = max(5.0, min(9.5, calc_font_size))
            
            # 折り返した行同士の隙間は少し詰める
            inner_line_spacing = ppt_font_size * 1.15
            
            # テキスト自体の総高さ
            total_text_height = num_lines * inner_line_spacing
            
            # 余った隙間を割り出して、各項目の「間」に均等に振り分ける
            remaining_height = available_height_pt - total_text_height
            # 🌟 修正：項目が0件の時にゼロ割りエラー（ZeroDivisionError）になるのを防ぐ安全ガードを追加！
            space_between = remaining_height / len(items) if remaining_height > 0 and len(items) > 0 else 0
            
            for idx, item_text in enumerate(items):
                p = tf_info.paragraphs[0] if idx == 0 else tf_info.add_paragraph()
                p.text = item_text
                
                p.font.size = Pt(ppt_font_size) 
                p.font.name = "游明朝"
                p.font.color.rgb = RGBColor(80, 80, 80)
                p.line_spacing = Pt(inner_line_spacing)
                
                # 🌟【究極の整列魔法】パワポの裏側（XML）を直接操作してインデントを設定！
                # ※ 8.5の部分は文字幅です。もし「:」と文字の隙間が広すぎたり狭すぎる場合は、この数字を微調整してください
                indent_emu = int(Pt(ppt_font_size * 8.5)) 
                pPr = p._p.get_or_add_pPr()
                pPr.set('marL', str(indent_emu))      # 2行目以降の開始位置(左余白)を設定
                pPr.set('indent', str(-indent_emu))   # 1行目(項目名)だけ左に引き戻す
                
                # 🌟 等間隔に散らすために、項目と項目の間に計算した余白を挿入！
                if idx > 0:
                    p.space_before = Pt(space_between)

            line_eq1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(5.9), Inches(10), Pt(1))
            line_eq1.fill.solid()
            line_eq1.fill.fore_color.rgb = gray_line
            line_eq1.line.fill.background()

            tb_eq = slide.shapes.add_textbox(Inches(0.2), Inches(6.05), Inches(1.0), Inches(0.4))
            p_eq = tb_eq.text_frame.paragraphs[0]
            p_eq.text = "物件設備"
            p_eq.font.size = Pt(11)
            p_eq.font.bold = True

            # 🌟 ここから変更：アイコン画像があれば貼り付ける処理
            icons = [icon_image1, icon_image2, icon_image3, icon_image4, icon_image5, icon_image6]
            import tempfile
            
            for i in range(6):
                current_left = Inches(1.2 + i * 1.4)
                current_top = Inches(6.05)
                icon_file = icons[i]
                
                # 画像がアップロードされているかチェック
                if icon_file and icon_file.filename:
                    suffix = os.path.splitext(icon_file.filename)[1]
                    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
                        tmp.write(icon_file.file.read()) # 画像データを書き込む
                        tmp_path = tmp.name
                    
                    # 🌟 変更：サイズを大きくし、縦横比を保ったまま中央に配置する
                    from PIL import Image
                    with Image.open(tmp_path) as img:
                        img_w, img_h = img.size
                    
                    # 🌟 アイコンの高さを0.6インチ（元の枠0.4インチの1.5倍）に拡大！
                    target_h_inches = 0.6
                    target_w_inches = target_h_inches * (img_w / img_h)
                    
                    # 枠の横幅(1.2インチ)と高さ(0.4インチ)との差分から、中央寄せするためのXY座標を計算
                    offset_x_inches = (1.2 - target_w_inches) / 2
                    offset_y_inches = (0.4 - target_h_inches) / 2
                    
                    centered_left = current_left + Inches(offset_x_inches)
                    centered_top = current_top + Inches(offset_y_inches)
                    
                    # 計算した中央の座標に、大きくなった画像を配置
                    slide.shapes.add_picture(tmp_path, centered_left, centered_top, height=Inches(target_h_inches))
                    os.remove(tmp_path)
                else:
                    # 🌟 修正：アイコン画像がない場合も、背景に合わせて真っ白な枠を出す
                    add_color_box(current_left, current_top, Inches(1.2), Inches(0.4), "アイコン", RGBColor(255, 255, 255), 11)

            line_eq2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.6), Inches(10), Pt(2))
            line_eq2.fill.solid()
            line_eq2.fill.fore_color.rgb = RGBColor(30, 45, 61)
            line_eq2.line.fill.background()

# ==========================================
        # 🎨 共通フッターエリア（スマート・モダンフォント完全対応版！）
        # ==========================================
        footer_y = Inches(6.45)
        footer_h = Inches(1.05)

        # 🌟 候補2：ラグジュアリー・クラシック用のフォント定義
        FONT_ENG = "Trebuchet MS"
        FONT_JPN = "メイリオ (Meiryo)"

        # 🌟 背景と基本色の設定（デフォルトはTOHOブルー）
        bg_color = RGBColor(10, 50, 96) 
        text_color_main = RGBColor(255, 255, 255)
        text_color_sub = RGBColor(200, 200, 200)
        line_color = RGBColor(255, 255, 255)
        top_border_color = RGBColor(10, 50, 96)
        invert_logo = True

        if footer_design_num == "2": # 白背景
            bg_color = RGBColor(255, 255, 255)
            text_color_main = RGBColor(51, 51, 51)
            text_color_sub = RGBColor(119, 119, 119)
            line_color = RGBColor(207, 216, 220)
            top_border_color = RGBColor(10, 50, 96)
            invert_logo = False
        elif footer_design_num == "3": # グレー背景
            bg_color = RGBColor(244, 246, 248)
            text_color_main = RGBColor(51, 51, 51)
            text_color_sub = RGBColor(119, 119, 119)
            line_color = None
            top_border_color = None
            invert_logo = False

        # --- 背景の描画 ---
        if footer_design_num == "4" and custom_footer_image:
            # 🌟 オリジナル画像を使用
            import tempfile
            from PIL import Image
            suffix = os.path.splitext(custom_footer_image.filename)[1]
            if not suffix: suffix = ".png"
            with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
                tmp.write(custom_footer_image.file.read())
                tmp_path = tmp.name
            slide.shapes.add_picture(tmp_path, 0, footer_y, width=Inches(10), height=footer_h)
            os.remove(tmp_path)
        else:
            # 帯の背景
            bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, footer_y, Inches(10), footer_h)
            bg_shape.fill.solid()
            bg_shape.fill.fore_color.rgb = bg_color
            bg_shape.line.fill.background()
            
            # 上のボーダー線
            if top_border_color:
                top_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, footer_y, Inches(10), Pt(2))
                top_line.fill.solid()
                top_line.fill.fore_color.rgb = top_border_color
                top_line.line.fill.background()

            # ロゴ
            base_dir = os.path.dirname(__file__)
            logo_path = os.path.join(base_dir, "static", "logo.png")
            if branch_name == "練馬": logo_path = os.path.join(base_dir, "static", "logo_nerima.png")
            elif branch_name == "武蔵野": logo_path = os.path.join(base_dir, "static", "logo_musashino.jpg")

            if os.path.exists(logo_path):
                # TOHOブルーの時はロゴを白抜きに変換
                if invert_logo and logo_path.endswith('.png'):
                    from PIL import Image
                    from io import BytesIO
                    with Image.open(logo_path).convert("RGBA") as img:
                        r, g, b, a = img.split()
                        r = r.point(lambda i: 255)
                        g = g.point(lambda i: 255)
                        b = b.point(lambda i: 255)
                        white_img = Image.merge("RGBA", (r, g, b, a))
                        tmp_logo_io = BytesIO()
                        white_img.save(tmp_logo_io, format="PNG")
                        tmp_logo_io.seek(0)
                        # 🌟 左に引っ張って透明な余白を相殺し、左右の中央に配置
                        slide.shapes.add_picture(tmp_logo_io, Inches(-0.05), footer_y - Inches(0.075), height=Inches(1.2))
                else:
                    # 🌟 こちらも同様に修正
                    slide.shapes.add_picture(logo_path, Inches(-0.05), footer_y - Inches(0.075), height=Inches(1.2))

            # 縦線1
            if line_color:
                vl1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.6), footer_y + Inches(0.15), Pt(1), Inches(0.75))
                vl1.fill.solid()
                vl1.fill.fore_color.rgb = line_color
                vl1.line.fill.background()

            # 🌟 絶対に綺麗に出るフォント定義
            FONT_ENG = "Arial"
            FONT_JPN = "游明朝"

            # 会社情報（🌟太字を減らして抜け感を出す）
            tb_comp = slide.shapes.add_textbox(Inches(2.7), footer_y + Inches(0.08), Inches(3.5), Inches(0.85))
            tf_comp = tb_comp.text_frame
            tf_comp.clear()
            tf_comp.margin_top = 0
            
            p1 = tf_comp.paragraphs[0]
            r1 = p1.add_run()
            r1.text = "TEL."
            r1.font.size = Pt(11)
            r1.font.bold = False # 🌟 細字にしてスタイリッシュに
            r1.font.color.rgb = text_color_sub
            r1.font.name = FONT_ENG
            
            r2 = p1.add_run()
            r2.text = f" {branch.get('tel', '')}"
            r2.font.size = Pt(24)
            r2.font.bold = False # 🌟 あえて細字にして高級感を出す
            r2.font.color.rgb = text_color_main
            r2.font.name = FONT_ENG
            
            p2 = tf_comp.add_paragraph()
            p2.text = branch.get("license", "免許番号")
            p2.font.size = Pt(8)
            p2.font.bold = False
            p2.font.color.rgb = text_color_sub
            p2.font.name = FONT_JPN
            
            p3 = tf_comp.add_paragraph()
            p3.text = branch.get('full_name', '')
            p3.font.size = Pt(13)
            p3.font.bold = True # 会社名だけ太字
            p3.font.color.rgb = text_color_main
            p3.font.name = FONT_JPN
            
            p4 = tf_comp.add_paragraph()
            p4.text = branch.get('address', '')
            p4.font.size = Pt(8)
            p4.font.bold = False
            p4.font.color.rgb = text_color_sub
            p4.font.name = FONT_JPN

            # 縦線2
            if line_color:
                vl2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.0), footer_y + Inches(0.15), Pt(1), Inches(0.75))
                vl2.fill.solid()
                vl2.fill.fore_color.rgb = line_color
                vl2.line.fill.background()

            # 担当者情報（🌟美しく整列）
            tb_person = slide.shapes.add_textbox(Inches(6.1), footer_y + Inches(0.15), Inches(2.0), Inches(0.75))
            tf_person = tb_person.text_frame
            tf_person.clear()
            tf_person.margin_top = 0
            
            p_p1 = tf_person.paragraphs[0] 
            r1 = p_p1.add_run()
            r1.text = "担当 "
            r1.font.size = Pt(9)
            r1.font.bold = False
            r1.font.color.rgb = text_color_sub
            r1.font.name = FONT_JPN
            
            r2 = p_p1.add_run()
            r2.text = "担当者名"
            r2.font.size = Pt(14)
            r2.font.bold = True
            r2.font.color.rgb = text_color_main
            r2.font.name = FONT_JPN
            
            p_p2 = tf_person.add_paragraph()
            p_p2.space_before = Pt(4)
            r_mob_lbl = p_p2.add_run()
            r_mob_lbl.text = "Mobile. "
            r_mob_lbl.font.size = Pt(8.5)
            r_mob_lbl.font.bold = False
            r_mob_lbl.font.color.rgb = text_color_main
            r_mob_lbl.font.name = FONT_ENG
            
            r_mob_val = p_p2.add_run()
            r_mob_val.text = "携帯番号"
            r_mob_val.font.size = Pt(8.5)
            r_mob_val.font.bold = False
            r_mob_val.font.color.rgb = text_color_sub
            r_mob_val.font.name = FONT_ENG # 数字なので英字フォント
            
            p_p3 = tf_person.add_paragraph()
            p_p3.space_before = Pt(1)
            r_mail_lbl = p_p3.add_run()
            r_mail_lbl.text = "Mail. "
            r_mail_lbl.font.size = Pt(8.5)
            r_mail_lbl.font.bold = False
            r_mail_lbl.font.color.rgb = text_color_main
            r_mail_lbl.font.name = FONT_ENG
            
            r_mail_val = p_p3.add_run()
            r_mail_val.text = "メールアドレス"
            r_mail_val.font.size = Pt(8.5)
            r_mail_val.font.bold = False
            r_mail_val.font.color.rgb = text_color_sub
            r_mail_val.font.name = FONT_ENG

            # 縦線3
            if line_color:
                vl3 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.1), footer_y + Inches(0.15), Pt(1), Inches(0.75))
                vl3.fill.solid()
                vl3.fill.fore_color.rgb = line_color
                vl3.line.fill.background()

            # 取引態様・手数料の表
            table_left = Inches(8.3)
            table_top = footer_y + Inches(0.2)
            col_w = Inches(0.8)
            row_h1 = Inches(0.25)
            row_h2 = Inches(0.35)
            
            def draw_cell(l, t, w, h, text, is_header):
                rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
                rect.fill.solid()
                if footer_design_num == "1":
                    # 🌟 修正：上段は青、下段は白（RGB(255,255,255)）にする！
                    rect.fill.fore_color.rgb = bg_color if is_header else RGBColor(255, 255, 255)
                    # 下段の枠線は白背景でも見えやすいように薄いグレーにする
                    rect.line.color.rgb = RGBColor(255, 255, 255) if is_header else RGBColor(200, 200, 200)
                    tc = RGBColor(255, 255, 255) if is_header else RGBColor(0, 0, 0)
                elif footer_design_num == "2":
                    rect.fill.fore_color.rgb = top_border_color if is_header else RGBColor(255,255,255)
                    rect.line.color.rgb = top_border_color
                    tc = RGBColor(255,255,255) if is_header else RGBColor(0,0,0)
                else:
                    rect.fill.fore_color.rgb = RGBColor(136,136,136) if is_header else RGBColor(255,255,255)
                    rect.line.color.rgb = RGBColor(136,136,136)
                    tc = RGBColor(255,255,255) if is_header else RGBColor(0,0,0)

                rect.line.width = Pt(1)
                tf = rect.text_frame
                tf.word_wrap = True
                tf.margin_top = tf.margin_bottom = tf.margin_left = tf.margin_right = 0
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                p = tf.paragraphs[0]
                p.text = text
                p.font.size = Pt(9)
                p.font.color.rgb = tc
                p.font.name = FONT_JPN # 表の中も美しく統一
                p.alignment = PP_ALIGN.CENTER

            draw_cell(table_left, table_top, col_w, row_h1, "取引態様", True)
            draw_cell(table_left + col_w, table_top, col_w, row_h1, "手数料", True)
            draw_cell(table_left, table_top + row_h1, col_w, row_h2, "", False)
            draw_cell(table_left + col_w, table_top + row_h1, col_w, row_h2, "", False)

        # 6. 完成したPowerPointファイルをバイナリ化して返す
        from fastapi.responses import StreamingResponse
        from io import BytesIO
        pptx_io = BytesIO()
        prs.save(pptx_io)
        pptx_io.seek(0)
        
        headers = {
            'Content-Disposition': f'attachment; filename="zumen_design{design_num}.pptx"'
        }
        return StreamingResponse(
            pptx_io, 
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers=headers
        )

    except Exception as e:
        import traceback
        from fastapi.responses import JSONResponse
        error_details = traceback.format_exc()
        print(f"Error: {error_details}")
        return JSONResponse(status_code=500, content={"detail": str(e)})

@app.post("/apb/extract_address")
async def apb_extract_address(zumen_file: UploadFile = File(...)):
    try:
        if not gemini_client:
            return {"status": "error", "message": "Geminiクライアントが準備されていません"}
        
        zumen_bytes = await zumen_file.read()
        mime_type = "application/pdf" if zumen_file.filename.lower().endswith(".pdf") else "image/jpeg"
        
        # 住所の特定に特化した超軽量プロンプト
        ocr_analysis = gemini_client.models.generate_content(
            model='gemini-2.5-flash',
            contents=[
                "この図面から『物件の正確な所在地（住所）』を特定し、その住所の文字列だけを出力してください。余計な解説や文字は一切含めず、住所のみ（例：東京都西東京市谷戸町3-28-16）にしてください。",
                types.Part.from_bytes(data=zumen_bytes, mime_type=mime_type)
            ]
        )
        
        address = ocr_analysis.text.strip() if ocr_analysis else ""
        print(f"📍 AIが特定した住所: {address}")
        return {"status": "success", "address": address}
        
    except Exception as e:
        return {"status": "error", "message": str(e)}