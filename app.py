import streamlit as st
import os
import json
import fitz  # PyMuPDF
import glob
import time
import random
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from google import genai
from dotenv import load_dotenv
from io import BytesIO
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageDraw, ImageFont, ImageEnhance
from google.genai import types

# --- 1. 環境設定 ---
load_dotenv()
gemini_api_key = os.getenv("GEMINI_API_KEY")

st.set_page_config(page_title="AI Pamphlet Builder", layout="wide")
st.title("🏡 AI Pamphlet Builder (Template Mode)")
st.write("用意したテンプレート画像とAIのテキスト解析を組み合わせ、全6ページのパワポ資料を安定して自動生成します。")

# フォントのパス（実行環境に合わせて確認してください）
FONT_PATH = './NotoSansCJKjp-Bold.ttf' 

# --- 修正：デザインテーマを6種＋その他に拡張 ---
THEMES = {
    "luxury":       {"name": "1 高級・ラグジュアリー", "bg_color": (40, 40, 45), "text_color": "white", "accent_color": (180, 150, 80)},
    "family":       {"name": "2 ファミリー・温もり", "bg_color": (255, 245, 235), "text_color": "black", "accent_color": (240, 130, 50)},
    "modern":       {"name": "3 スタイリッシュ・モダン", "bg_color": (240, 245, 255), "text_color": "black", "accent_color": (50, 100, 180)},
    "wa_modern":    {"name": "4 和モダン・伝統美", "bg_color": (230, 225, 215), "text_color": "black", "accent_color": (100, 120, 80)},
    "casual":       {"name": "5 カジュアル・ポップ", "bg_color": (255, 250, 220), "text_color": "black", "accent_color": (250, 100, 130)},
    "other":        {"name": "6 その他（自由入力スタイル）", "bg_color": (240, 240, 240), "text_color": "black", "accent_color": (100, 100, 100)}
}
# --- 追加：店舗ごとの詳細データ ---
BRANCH_DATA = {
    "練馬": {
        "full_name": "株式会社 東宝ハウス練馬",
        "license": "東京都知事（4）第86488号",
        "address": "〒178-0063 東京都練馬区東大泉1-27-22光和ビル2F",
        "tel": "0120-384-700"
    },
    "国分寺": {
        "full_name": "株式会社 東宝ハウス国分寺",
        "license": "東京都知事（9）第42787号",
        "address": "〒185-0021 東京都国分寺市南町3-22-2",
        "tel": "0120-13-3107"
    },
    "武蔵野": {
        "full_name": "株式会社 東宝ハウス武蔵野",
        "license": "東京都知事（3）第90333号",
        "address": "〒180-0004 東京都武蔵野市吉祥寺本町1-15-9",
        "tel": "0120-15-3101"
    },
}

# --- 2. メモリ（セッション状態）の初期化 ---
if "finished_pages" not in st.session_state:
    st.session_state.finished_pages = []
if "ai_data" not in st.session_state:
    st.session_state.ai_data = None
if "current_file" not in st.session_state:
    st.session_state.current_file = None
if "pdf_text" not in st.session_state:
    st.session_state.pdf_text = ""
if "cover_choices" not in st.session_state:
    st.session_state.cover_choices = []
if "selected_cover_index" not in st.session_state:
    st.session_state.selected_cover_index = 0    
    

# --- 修正：空室写真アップロード欄を追加し、レイアウト選択UIを削除 ---
col_u1, col_u2 = st.columns(2)
with col_u1:
    uploaded_file = st.file_uploader("販売図面（PDFまたは画像）をアップロード", type=["pdf", "png", "jpg", "jpeg"])
with col_u2:
    madori_file = st.file_uploader("間取り図の画像をアップロード（P.4用）", type=["png", "jpg", "jpeg"])
    # ✨ 追加：空室写真のアップロード
    empty_room_file = st.file_uploader("空室写真をアップロード（P.6 バーチャルステージング用）", type=["png", "jpg", "jpeg"])

st.write("---")
st.subheader("🗺️ 地図画像のアップロード")
map_file = st.file_uploader("地図に使用するマップ画像をアップロード（P.3 MAP & ACCESS用）", type=["png", "jpg", "jpeg"])

# ファイルリセット処理
if uploaded_file is not None and uploaded_file.name != st.session_state.current_file:
    st.session_state.finished_pages = []
    st.session_state.ai_data = None
    st.session_state.current_file = uploaded_file.name
    st.session_state.pdf_text = ""
# ✨ 修正：UIは削除し、内部処理用のデフォルト値を設定
selected_layout_key = "horizontal"

# --- 物件設定（画面からは非表示にし、裏側で「戸建て」に固定） ---
selected_property_key = "house"
selected_property_type_label = "戸建て"
selected_apt_scale = "house"

# 3. デザインテーマを「高級・ラグジュアリー」に固定
selected_style_key = "luxury"
theme_info = THEMES[selected_style_key]
custom_style_description = ""

# 4. 担当店舗の選択（ここで定義される）
st.write("---")
st.subheader("🏢 担当店舗の選択")
selected_branch_name = st.selectbox("担当店舗を選んでください：", list(BRANCH_DATA.keys()))

room_features_input = ""

# 6. スライドの向き選択
st.write("---")
st.subheader("📏 スライドの向きを選択")
orientation = st.radio("作成するパンフレットの向き：", ["横向き (Landscape)", "縦向き (Portrait)"], index=1)

# 👇👇👇 ここから追加（ページ選択UI） 👇👇👇
st.write("---")
st.subheader("🛠️ デバッグ・生成ページ設定")
st.caption("必要なページだけをチェックすると、時間とAPIコストを節約できます。")

page_options = {
    "cover": "P.1-2 表紙 (Imagen)",
    "aerial_map": "P.3 コンセプト (Placeholder)",
    "access": "P.4 地図 (Image Processing)",
    "floor_plan": "P.5 間取り (Imagen)",
    "interior_hq": "P.6 内観生成 (Imagen)",
    "interior": "P.7 ギャラリー (Layout)",
    "company": "P.8 会社案内 (Layout)"
}

selected_pages_keys = st.multiselect(
    "生成するページを選択：",
    options=list(page_options.keys()),
    default=list(page_options.keys()), # デフォルトは全選択
    format_func=lambda x: page_options[x]
)
# 👆👆👆 ここまで追加 👆👆👆

# 7. 生成ボタン（一括生成ボタン）
st.write("---")
btn_generate = st.button("🚀 パンフレットを生成する", disabled=not uploaded_file)

# 選択された店舗の詳細を取得
branch_info = BRANCH_DATA[selected_branch_name]

# --- 4. メイン処理（パンフレット一括生成） ---
gemini_client = genai.Client(api_key=gemini_api_key)

def generate_with_retry(model_name, prompt_contents, generation_config=None):
    # ✨ 修正1：待機時間を細かく、試行回数を増やしました
    wait_times = [30, 60, 120, 180, 240, 300, 360, 420, 480, 540, 600]
    max_attempts = len(wait_times) + 1
    
    # ✨ 修正2：メッセージが下に伸びないよう、専用の「入れ替え用の枠」を作ります
    warning_box = st.empty() 
    
    for attempt in range(max_attempts):
        try:
            time.sleep(2)
            if generation_config:
                res = gemini_client.models.generate_content(model=model_name, contents=prompt_contents, config=generation_config)
            else:
                res = gemini_client.models.generate_content(model=model_name, contents=prompt_contents)
            
            # 無事に生成できたら、警告メッセージの枠を消して綺麗にする
            warning_box.empty() 
            return res
            
        except Exception as inner_e:
            if ("503" in str(inner_e) or "429" in str(inner_e)) and attempt < len(wait_times):
                wait_sec = wait_times[attempt]
                # st.warningではなく、先ほど作った枠(warning_box)の中身を書き換える
                warning_box.warning(f"サーバー制限に到達しました。{wait_sec}秒待機して再試行します... (再試行 {attempt + 1}/{len(wait_times)})")
                time.sleep(wait_sec)
            else:
                # 別の致命的なエラーの場合は枠を消してエラーを投げる
                warning_box.empty()
                raise inner_e

if btn_generate and uploaded_file is not None:
    with st.status("📄 パンフレットを作成中...（約2分かかります）", expanded=True) as status:
        progress_msg = st.empty()
        try:
            # --- 4-A. OCR ＆ 表紙生成 ---
            progress_msg.info("🔍 販売図面をAI（OCR）で読み取っています...")
            file_bytes = uploaded_file.getvalue()
            filename_lower = uploaded_file.name.lower()
            if filename_lower.endswith(".pdf"): mime_type = "application/pdf"
            elif filename_lower.endswith(".png"): mime_type = "image/png"
            else: mime_type = "image/jpeg"
            
            ocr_analysis = generate_with_retry(
                model_name='gemini-2.5-flash',
                prompt_contents=[
                    "図面に記載されているすべてのテキスト情報を正確に抽出してください。特に「物件名」「所在地」「交通（最寄り駅情報）」「価格」「面積」「周辺施設」の情報は絶対に漏れなく抽出してください。Imagen生成用の地名（市区町村）も特定してください。",
                    types.Part.from_bytes(data=file_bytes, mime_type=mime_type)
                ]
            )
            st.session_state.pdf_text = ocr_analysis.text

# ✨ 修正：表紙が選択されている時だけ画像生成を実行し、3パターンのプロンプトを設定
            if "cover" in selected_pages_keys:
                progress_msg.info("📸 表紙のデザイン案を3パターン生成中...")
                target_city_town = "Tokyo"
                if "所在地" in st.session_state.pdf_text:
                    target_city_town = st.session_state.pdf_text.split("所在地")[-1][:10]

                prompts = [
                    # 案1：公園で遊ぶ幸せな4人家族（現状維持） [cite: 17-20]
                    f"High-end luxury magazine photography. A happy Japanese family (father, mother, and two children) playing happily in a beautiful sunny park in {target_city_town}. Bright natural daylight, soft sunlight filtering through trees. Modern and sophisticated casual style. NO text.",
                    
                    # 案2：一軒家の室内で寛ぐ男性（✨外の景色を最小限にする指示を追加） [cite: 20-22]
                    f"High-end luxury lifestyle photography. A sophisticated Japanese man relaxing and drinking coffee in a modern, spacious luxury living room of a detached house in {target_city_town}. Focus on the premium interior design, close-up shot. Windows are softly blurred or showing very minimal outdoor scenery to emphasize the indoor atmosphere. NO text.",
                    
                    # 案3：マンションの室内で遊ぶ家族(✨外の景色を最小限にする指示を追加、タワマン感を排除） [cite: 22]
                    f"High-end luxury lifestyle photography. A happy Japanese family playing in a stylish, high-ceiling living room of a luxury modern apartment in {target_city_town}. Large windows are softly blurred or show very minimal outdoor scenery, carefully avoiding a distant tower apartment skyline view to emphasize the cozy indoor atmosphere. Sophisticated interior design, bright natural light, warm family atmosphere. NO text."
                ]

                temp_choices = []
                for p in prompts:
                    img_res = gemini_client.models.generate_images(
                        model='imagen-4.0-generate-001',
                        prompt=p,
                        config=types.GenerateImagesConfig(
                            number_of_images=1, 
                            aspect_ratio="4:3" if orientation == "横向き (Landscape)" else "3:4"
                        )
                    )
                    temp_choices.append(Image.open(BytesIO(img_res.generated_images[0].image.image_bytes)))
                st.session_state.cover_choices = temp_choices
            else:
                st.write("⏭️ 表紙が未選択のため、画像生成をスキップします。")
            # --- 4-B. 全ページデータ（JSON）の生成 ---

# --- 4-B. 全ページデータ（JSON）の生成 ---
            progress_msg.info("🔍 AIが物件情報、指定地域、デザインテーマを分析中...")
            
            theme_info = THEMES[selected_style_key]
            current_theme_name = custom_style_description if selected_style_key == "other" else theme_info["name"]

            # マンションの外観サンプル画像を解析
            apt_style_prompt = ""
            if selected_property_key == "apartment":
                progress_msg.info(f"🏢 マンション規模（{selected_apt_scale}）のサンプル画像を解析中...")
                sample_img_path = f"apt_{selected_apt_scale}.jpg" 
                
                if os.path.exists(sample_img_path):
                    with open(sample_img_path, "rb") as f:
                        s_bytes = f.read()
                    try:
                        analysis_apt = generate_with_retry(
                            model_name='gemini-2.5-flash',
                            prompt_contents=[
                                "Analyze this apartment building's exterior design (colors, materials, window styles, overall architectural shape). Generate a short, STRICT English prompt. Output ONLY the English prompt string.",
                                types.Part.from_bytes(data=s_bytes, mime_type="image/jpeg")
                            ]
                        )
                        apt_style_prompt = analysis_apt.text.strip()
                    except Exception as e:
                        st.warning(f"マンション画像の解析エラー: {e}")
                else:
                    progress_msg.info(f"※ サンプル画像 ({sample_img_path}) が見つからないため、通常生成します。")

            # 間取り図解析
            room_description = "A modern living room"
            if madori_file:
                progress_msg.info("🔍 間取り図から強力な画像生成プロンプトを作成中...")
                m_bytes = madori_file.getvalue()
                
                # ✨ 追加：選択されたレイアウト構図の英語指示を定義
                layout_instruction = ""
                if selected_layout_key == "horizontal":
                    layout_instruction = "CAMERA ANGLE: Wide shot, landscape composition, parallel to the main window. Emphasize the horizontal breadth and spaciousness of the room. Bright, open feeling with massive windows spanning the background."
                else:
                    layout_instruction = "CAMERA ANGLE: Deep shot, one-point perspective, looking down the length of the room (e.g., from the dining area in the foreground towards the living area and window in the background). Emphasize depth and structured spatial arrangement."

                user_instruction = f"\n\n[USER'S ABSOLUTE REQUIREMENTS]\n{room_features_input}\n\n[COMPOSITION REQUIREMENT]\n{layout_instruction}"
                
                # ✨ 便利関数を使って間取り図を解析
                analysis = generate_with_retry(
                    model_name='gemini-2.5-flash',
                    prompt_contents=[
                        f"""Analyze the attached Japanese floor plan (LDK area) and generate a highly detailed, STRICT English prompt for an Image Generation AI.

{user_instruction}

[CRITICAL RULES TO OVERCOME AI BIAS]
1. Kitchen (Crucial): Image AIs always default to island kitchens. If the floor plan or user indicates a wall-mounted kitchen (壁付けキッチン), you MUST explicitly write: "single-wall kitchen, cabinets and stove placed flat against the back wall, NO island, NO peninsula, wide open floor space in front".
2. Windows: If a large window is requested, explicitly write: "massive wall-to-wall and floor-to-ceiling panoramic windows, clear view, no vertical pillars blocking the view".
3. Layout: Specify exactly where things are based on the floor plan (e.g., "kitchen on the left side, dining table in the center, large window at the far back"). Ensure it aligns perfectly with the COMPOSITION REQUIREMENT.
4. Style: Match the theme "{current_theme_name}". High-end architectural photography, 8k resolution.

[OUTPUT FORMAT]
Output ONLY the final English prompt string. Do NOT output any conversational text like "Here is the prompt". Do NOT use markdown blocks. Start directly with the description.""",
                        types.Part.from_bytes(data=m_bytes, mime_type="image/jpeg")
                    ]
                )
                room_description = analysis.text.strip()
                
            elif room_features_input:
                # 間取り図がなく、テキスト入力だけがある場合
                trans = generate_with_retry(
                    model_name='gemini-2.5-flash',
                    prompt_contents=[f"Translate the following interior design requirements into a highly detailed English prompt for an image generation AI. Output ONLY the English prompt string, no conversational text.\nRequirements: {room_features_input}\nTheme: {current_theme_name}"]
                )
                room_description = trans.text.strip()
            
            ratio_text = "4:3（横長）" if orientation == "横向き (Landscape)" else "3:4（縦長）"
            
            prompt = f"""
            あなたはプロの不動産ライターです。提供された【補助データ】を隅々まで解析し、以下の7ページ構成のパンフレット（比率 {ratio_text}）を作成してください。
            
            【データ抽出の絶対条件】
            1. **property_name_en**: 【補助データ】から物件名を特定してください。雑誌風のデザインにするため、可能な限り英語表記にするか、読みをアルファベットに変換してください。（例：ラ・フレーズ国分寺 → La Phrase Kokubunji）
            2. **city_town**: 物件の「所在地」から市区町村と町名を抽出してください。（例：東京都府中市美好町... → 府中市 美好町）
            3. **station_info**: 「交通」の項目から、最も主要な駅名と徒歩分数を抽出してください。改行を入れて見やすくしてください。（例：中央線 国分寺駅 徒歩5分 → 国分寺駅\n徒歩5分）
            4. **price**: 販売価格を抽出してください。
            5. **land_area**: 【補助データ】から土地面積を抽出してください。「土地」という文字と「㎡」単位を含めてください。（例：土地 100.77㎡）
            6. **building_area**: 【補助データ】から建物面積を抽出してください。「建物」という文字と「㎡」単位を含めてください。（例：建物 100.77㎡）
            7. **price_jp**: 【補助データ】から価格を日本語表記（万円単位）で抽出してください。（例：3,650万円）
            8. **デザインスタイル**: {current_theme_name} の雰囲気に合わせた魅力的な言葉選びをしてください。
            9. **property_name_jp**: 【補助データ】のOCRテキストの中から、最も目立つ場所に書かれている名称（通常はタイトル）、または物件概要欄から特定される物件名を日本語表記で特定し、抽出してください。部屋番号（〜号室など）の記載がある場合は必ず含めてください。（例：パークシティひばりが丘 410号室）
            10. **city_town_clean**: 【補助データ】の所在地から市区町村と町名を抽出し、余計な文字（東京都など）を省いてください（例：府中市美好町）。これをImagen生成用に保持します。

            出力は必ず以下のJSON配列形式のみにしてください。
            [
              {{
                "page": 1, "type": "cover",
                "price": "抽出した価格", 
                "property_name_en": "抽出・変換した物件名（英字）", 
                "property_name_jp": "抽出した物件名（日本語）",
                "sub_copy": "テーマに合わせた洗練されたサブコピー（日本語）",
                "city_town": "抽出した地域名（例：府中市 美好町）",
                "station_info": "抽出した駅名と徒歩分数（例：国分寺駅\\n徒歩5分）"
              }},
              {{
                "page": 2, "type": "aerial_map",
                "headline": "FUTURE VISION", 
                "sub_headline": "未来を描く",
                "main_text": "【補助データ】の周辺環境情報を基に、このテーマの客層に刺さる紹介文を3〜4行で作成してください。",
                "city_town_clean": "抽出したImagen用地域名（例：府中市美好町）"
              }},
              {{
                "page": 3, "type": "access", "headline": "MAP & ACCESS", 
                "source_pdf_page": 1, 
                "access_info": "【補助データ】から交通アクセス情報（路線・駅名・徒歩分数など）を箇条書きで抽出してください。",
                "life_info": "【補助データ】の「周辺施設」などの項目から、施設名と距離（徒歩〇分など）を箇条書きで漏れなく抽出してください。"
              }},
              {{
                "page": 4, "type": "floor_plan",
                "property_name_jp": "抽出した物件名（日本語）",
                "land_area": "抽出した土地面積（例：土地 --- ㎡）",
                "building_area": "抽出した建物面積（例：建物 --- ㎡）",
                "price_jp": "抽出した価格（日本語表記、例：--- 万円）"
              }},
              {{
                "page": 5, "type": "interior_hq", "headline": "INTERIOR VISION", 
                "sub_headline": "間取り図から描き出した完成予想イメージ"
              }},
              {{
                "page": 6, "type": "interior", "headline": "内観ギャラリー", "source_pdf_page": 1 
              }},
              {{ 
                "page": 7, "type": "company", 
                "company_name": "{branch_info['full_name']}",
                "license": "{branch_info['license']}",
                "address": "{branch_info['address']}",
                "tel": "{branch_info['tel']}"
              }}
            ]
            
            【補助データ（図面から抽出されたテキスト）】:
            {st.session_state.pdf_text}
            """
            
            # ✨ 便利関数を使ってJSON（文章）を生成
            response = generate_with_retry(
                model_name='gemini-2.5-flash', 
                prompt_contents=prompt, 
                generation_config=types.GenerateContentConfig(response_mime_type="application/json")
            )
            
            raw_json = response.text.strip()
            if raw_json.startswith("```json"): 
                raw_json = raw_json[7:-3].strip()
            elif raw_json.startswith("```"): 
                raw_json = raw_json[3:-3].strip()
                
            st.session_state.ai_data = json.loads(raw_json)
            
            # --- 描画用のヘルパー関数群 ---
            # ✨ 修正：PDFだけでなくJPG画像にも対応させた関数
            def get_source_image(file_bytes, filename, target_page_num):
                filename_lower = filename.lower()
                if filename_lower.endswith(".pdf"):
                    doc_temp = fitz.open(stream=file_bytes, filetype="pdf")
                    p_num = max(0, min(target_page_num, doc_temp.page_count - 1))
                    pix = doc_temp.load_page(p_num).get_pixmap(matrix=fitz.Matrix(2.5, 2.5))
                    return Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                else:
                    return Image.open(BytesIO(file_bytes)).convert("RGB")

            def get_fitting_font(draw_obj, text, initial_size, max_width):
                size = initial_size
                try: 
                    font = ImageFont.truetype(FONT_PATH, size)
                except: 
                    return ImageFont.load_default()
                while draw_obj.textbbox((0, 0), text, font=font)[2] > max_width and size > 10:
                    size -= 2
                    font = ImageFont.truetype(FONT_PATH, size)
                return font

            def draw_dashed_rectangle(draw, rect, outline="gray", width=2):
                x1, y1, x2, y2 = rect
                dash = 10
                for x in range(int(x1), int(x2), dash * 2): draw.line([(x, y1), (min(x + dash, x2), y1)], fill=outline, width=width)
                for x in range(int(x1), int(x2), dash * 2): draw.line([(x, y2), (min(x + dash, x2), y2)], fill=outline, width=width)
                for y in range(int(y1), int(y2), dash * 2): draw.line([(x1, y), (x1, min(y + dash, y2))], fill=outline, width=width)
                for y in range(int(y1), int(y2), dash * 2): draw.line([(x2, y), (x2, min(y + dash, y2))], fill=outline, width=width)

            # ✨修正：画像単体ではなく、ページデータとセット（タプル）で保存するリストに変更
            generated_pages_with_data = [] 
            
            # 向きに合わせてキャンバスサイズを入れ替える
            if orientation == "横向き (Landscape)":
                width, height = 1024, 768
            else:
                width, height = 768, 1024
            
            theme_bg = theme_info["bg_color"]
            theme_tc = theme_info["text_color"]
            theme_ac = theme_info["accent_color"]

            for page_data in st.session_state.ai_data:
                p_type = page_data['type']
                
                # ✨修正：選択されていないページは完全にスキップ（ダミー画像も入れない）
                if p_type not in selected_pages_keys:
                    st.write(f"⏭️ {page_data['page']}ページ目（{p_type}）はスキップします")
                    continue

                st.write(f"🖼️ {page_data['page']}ページ目（{p_type}）を合成中...")
                bg_image = None
                
# ────────── ① 表紙（✨生成されたすべてのパターンを追加するループに修正） ──────────
                if p_type == 'cover':
                    if st.session_state.cover_choices:
                        # 全パターンの画像をループで追加
                        for cover_img in st.session_state.cover_choices:
                            resized_cover = cover_img.copy().resize((width, height))
                            generated_pages_with_data.append((page_data, resized_cover))
                    else:
                        bg_image = Image.new('RGB', (width, height), color=theme_bg)
                        generated_pages_with_data.append((page_data, bg_image))
                    continue # 表紙の処理はここで完了
                
                # ────────── ② コンセプト（旧：空撮地図） ──────────
                elif page_data['type'] == 'aerial_map':
                    st.write(f"🎨 コンセプトページのレイアウトを作成中...")

                    # ✨ 修正：AIによる画像生成をスキップし、PowerPoint側でプレースホルダーになるように設定
                    st.session_state.concept_img = None

                    # ✨ プレビュー画面用の背景（スッキリ分割レイアウト）
                    bg_image = Image.new('RGB', (width, height), color=(245, 240, 225)) # 上品なベージュ背景
                    draw = ImageDraw.Draw(bg_image)
                    
                    bronze_color = (215, 185, 140)
                    gray_color = (235, 235, 235)
                    
                    if orientation == "横向き (Landscape)":
                        # メイン画像枠をグレーで塗る（右半分）
                        draw.rectangle([(int(width * 0.45), 0), (width, height)], fill=gray_color)
                        # 写真1, 2の枠
                        frames = [
                            (width*0.3, height*0.6, width*0.55, height*0.86),
                            (width*0.75, height*0.06, width*0.95, height*0.4)
                        ]
                        # 装飾直線
                        draw.line([(width*0.05, height*0.17), (width*0.45, height*0.17)], fill=bronze_color, width=2)
                        
                        # テキスト位置
                        head_y, sub_y, main_y = height*0.06, height*0.21, height*0.33
                    else:
                        # 縦向き：メイン画像枠をグレーで塗る（下半分）
                        draw.rectangle([(0, int(height * 0.45)), (width, height)], fill=gray_color)
                        # 写真1, 2の枠
                        frames = [
                            (width*0.6, height*0.35, width*0.93, height*0.55),
                            (width*0.06, height*0.75, width*0.39, height*0.95)
                        ]
                        # 装飾直線
                        draw.line([(width*0.05, height*0.17), (width*0.95, height*0.17)], fill=bronze_color, width=2)
                        
                        # テキスト位置
                        head_y, sub_y, main_y = height*0.06, height*0.21, height*0.33

                    # サブ写真枠をプレビューに描画
                    for f in frames:
                        draw.rectangle(f, fill=(245, 245, 245), outline=bronze_color, width=2)

                    # --- テキスト描画 ---
                    try:
                        font_head = ImageFont.truetype(FONT_PATH, int(height * 0.05))
                        font_subhead = ImageFont.truetype(FONT_PATH, int(height * 0.035))
                        font_main = ImageFont.truetype(FONT_PATH, int(height * 0.018))
                    except:
                        font_head = font_subhead = font_main = ImageFont.load_default()

                    headline = "THE CONCEPT"
                    sub_headline = page_data.get('sub_headline', '五感を満たす、静謐の邸宅。')
                    main_text = page_data.get('main_text', '')

                    draw.text((width*0.05, head_y), headline, font=font_head, fill="black", anchor="la")
                    draw.text((width*0.05, sub_y), sub_headline, font=font_subhead, fill="black", anchor="la")
                    draw.multiline_text((width*0.05, main_y), main_text, font=font_main, fill=(80, 80, 80), spacing=8)

                # ────────── ③ アクセス・地図 ──────────
                elif page_data['type'] == 'access':
                    st.write(f"🎨 MAP & ACCESSの背景画像をAIで生成中...")
                    try:
                        # Imagenによるエレガントな背景の生成
                        image_result = gemini_client.models.generate_images(
                            model='imagen-4.0-generate-001',
                            prompt="A simple, elegant and abstract background for a luxury real estate map and access page. Clean white and soft gold geometric accents, minimal, premium corporate design, no text, no logos.",
                            config=types.GenerateImagesConfig(number_of_images=1, aspect_ratio="4:3" if orientation == "横向き (Landscape)" else "3:4")
                        )
                        generated_bytes = image_result.generated_images[0].image.image_bytes
                        bg_image = Image.open(BytesIO(generated_bytes)).convert("RGB").resize((width, height))
                    except Exception as e:
                        st.warning(f"3ページ目の画像生成エラー: {e}")
                        bg_image = Image.new('RGB', (width, height), color=(250, 250, 250))

                    # 画面全体に「白い半透明のフィルター」を強めに被せて、上に乗る文字や地図を見やすくする
                    overlay = Image.new('RGBA', bg_image.size, (255, 255, 255, 220)) 
                    bg_image = Image.alpha_composite(bg_image.convert('RGBA'), overlay).convert('RGB')
                    
# ────────── ④ 間取り図 ──────────
                elif page_data['type'] == 'floor_plan':
                    st.write(f"🎨 間取り図ページの高級背景（木目グラデーション）をAIで生成中...")
                    
                    # ✨ 修正：AIへのプロンプトを「下に行くほど暗くなる木目」と「質感のあるヘッダー」に強化
                    bg_prompt = (
                        "A cinematic, high-end architectural background for a luxury property brochure. "
                        "Top 20% area: dark charcoal gray slate stone with elegant natural texture. "
                        "Bottom 80% area: rich, deep espresso luxury wood grain texture, vertical planks, "
                        "with a subtle vertical gradient that gets progressively darker towards the bottom edge. "
                        "High contrast, moody lighting, NO text, NO floor plans, NO logos."
                    )
                    
                    try:
                        image_result = gemini_client.models.generate_images(
                            model='imagen-4.0-generate-001',
                            prompt=bg_prompt,
                            config=types.GenerateImagesConfig(
                                number_of_images=1, 
                                aspect_ratio="4:3" if orientation == "横向き (Landscape)" else "3:4"
                            )
                        )
                        generated_bytes = image_result.generated_images[0].image.image_bytes
                        bg_image = Image.open(BytesIO(generated_bytes)).convert("RGB").resize((width, height))
                    except Exception as e:
                        st.warning(f"4ページ目の背景生成エラー: {e}")
                        bg_image = Image.new('RGB', (width, height), color=(25, 20, 15))

                    draw = ImageDraw.Draw(bg_image)
                    header_h = int(height * 0.20)

                    # ✨ 修正：ゴールドの区切り線をより金属感のある立体的な重なりに
                    gold_light = (215, 185, 140)
                    gold_dark = (160, 130, 80)
                    # 境界に太めのゴールド帯
                    draw.rectangle([(0, header_h - 15), (width, header_h + 15)], fill=gold_light)
                    # 中央に濃い色の細い線を入れて立体感を出す
                    draw.rectangle([(0, header_h - 2), (width, header_h + 2)], fill=gold_dark)

                    # 2. ユーザーがアップロードした正確な間取り図を配置
                    if madori_file:
                        m_img = Image.open(madori_file).convert("RGB")
                        frame_width = 15
                        m_img_with_frame = Image.new('RGB', (m_img.width + frame_width*2, m_img.height + frame_width*2), gold_light)
                        draw_frame = ImageDraw.Draw(m_img_with_frame)
                        draw_frame.rectangle([(5, 5), (m_img_with_frame.width-5, m_img_with_frame.height-5)], outline=(100, 80, 40), width=3)
                        
                        m_img_with_frame.paste(m_img, (frame_width, frame_width))
                        m_img_with_frame.thumbnail((width * 0.85, height * 0.70))
                        
                        p_x = (width - m_img_with_frame.width) // 2
                        p_y = header_h + 40 + (height - header_h - 40 - m_img_with_frame.height) // 2
                        bg_image.paste(m_img_with_frame, (int(p_x), int(p_y)))
                    else:
                        draw.text((width/2, height/2), "間取り図がアップロードされていません", fill="white", anchor="mm")
                
                # ────────── ⑤ 内観（✨空室写真ベースのバーチャルステージング対応） ──────────
                elif page_data['type'] == 'interior_hq':
                    st.write(f"🎨 完成予想イメージを生成中...")
                    
                    # ✨ 空室写真がアップロードされている場合は、その構造を解析して家具を合成する
                    if empty_room_file:
                        st.write("🔍 アップロードされた空室写真を解析し、構造を維持しつつ家具を配置します...")
                        e_bytes = empty_room_file.getvalue()
                        analysis_img = generate_with_retry(
                            model_name='gemini-2.5-flash',
                            prompt_contents=[
                                f"Analyze this empty room photo precisely. Generate a highly detailed English prompt for an Image Generation AI to recreate this EXACT room structure (same window placement, floor, wall) BUT add elegant {current_theme_name} style furniture. Output ONLY the prompt string.",
                                types.Part.from_bytes(data=e_bytes, mime_type="image/jpeg")
                            ]
                        )
                        final_room_prompt = analysis_img.text.strip()
                    else:
                        final_room_prompt = room_description # 空室写真がない場合は間取り図から

                    try:
                        image_result = gemini_client.models.generate_images(
                            model='imagen-4.0-generate-001',
                            prompt=final_room_prompt,
                            config=types.GenerateImagesConfig(number_of_images=1, aspect_ratio="4:3" if orientation == "横向き (Landscape)" else "3:4")
                        )
                        generated_bytes = image_result.generated_images[0].image.image_bytes
                        bg_image = Image.open(BytesIO(generated_bytes)).convert("RGB").resize((width, height))
                    except Exception as e:
                        st.error(f"内観生成エラー: {e}")
                        bg_image = Image.new('RGB', (width, height), color=(240, 240, 240))

                    # ✨ 画像の下部にシックな半透明帯を配置し、高級感を演出
                    overlay = Image.new('RGBA', bg_image.size, (0, 0, 0, 0))
                    draw_ov = ImageDraw.Draw(overlay)
                    draw_ov.rectangle([(0, height * 0.82), (width, height)], fill=(0, 0, 0, 180))
                    draw_ov.line([(0, height * 0.82), (width, height * 0.82)], fill=theme_ac + (255,), width=4)
                    bg_image = Image.alpha_composite(bg_image.convert('RGBA'), overlay).convert('RGB')
                    
                    draw = ImageDraw.Draw(bg_image)
                    f_h = get_fitting_font(draw, "INTERIOR VISION", int(height * 0.07), width * 0.5)
                    sub_text = page_data.get('sub_headline', '')
                    f_s = get_fitting_font(draw, sub_text, int(height * 0.03), width * 0.4)
                    draw.text((width*0.05, height*0.87), "INTERIOR VISION", font=f_h, fill="white", anchor="la")
                    draw.text((width*0.95, height*0.90), sub_text, font=f_s, fill="white", anchor="ra")
                
# ────────── ⑥ 内観ギャラリー ＆ ⑦ 会社案内 背景生成 ──────────
                elif page_data['type'] in ['interior', 'company']:
                    # ✨ 修正：P.7とP.8で背景を共通化し、統一感を出す
                    if "shared_luxury_bg" in st.session_state and st.session_state.shared_luxury_bg:
                        progress_msg.info(f"🖼️ 高級背景（共通デザイン）を適用中...")
                        bg_image = st.session_state.shared_luxury_bg.copy()
                    else:
                        progress_msg.info(f"🎨 高級背景（石目・スポットライト・金粉）をImagenで生成中...")
                        bg_prompt = (
                            "A cinematic, high-end luxury background for an interior gallery. "
                            "Deep charcoal black stone texture with a soft spotlight effect in the center. "
                            "Elegant golden dust particles. NO text, NO logos."
                        )
                        try:
                            image_result = gemini_client.models.generate_images(
                                model='imagen-3.0-generate-001', # ✨ 適切なモデル名に修正
                                prompt=bg_prompt,
                                config=types.GenerateImagesConfig(
                                    number_of_images=1, 
                                    aspect_ratio="4:3" if orientation == "横向き (Landscape)" else "3:4"
                                )
                            )
                            generated_bytes = image_result.generated_images[0].image.image_bytes
                            bg_image = Image.open(BytesIO(generated_bytes)).convert("RGB").resize((width, height))
                            st.session_state.shared_luxury_bg = bg_image.copy() # 保存
                        except Exception as e:
                            st.warning(f"背景生成エラー（予備デザインを適用）: {e}")
                            # ✨ 修正：AIが失敗しても真っ黒にせず、高級グラデーションを自作
                            bg_image = Image.new('RGB', (width, height), color=(20, 20, 25))
                            draw_grad = ImageDraw.Draw(bg_image)
                            for r in range(height, 0, -5):
                                alpha = int(100 * (r / height))
                                draw_grad.ellipse(
                                    [(width/2 - r, height/2 - r), (width/2 + r, height/2 + r)],
                                    outline=(30 + (100-alpha)//2, 30 + (80-alpha)//2, 40)
                                )
                            st.session_state.shared_luxury_bg = bg_image.copy()

                if bg_image:
                    generated_pages_with_data.append((page_data, bg_image))

            st.session_state.finished_pages = generated_pages_with_data
            status.update(label="✅ 全ページが完成しました！", state="complete", expanded=False)

        except Exception as e:
            st.error(f"作成中にエラーが発生しました: {e}")
            st.stop()

# --- 5. 画面表示とPowerPoint作成 ---
if st.session_state.finished_pages:
    # ✨ 追加：PowerPoint作成時にも、選んだテーマの文字色を認識させる
    theme_info = THEMES[selected_style_key]
    theme_tc = theme_info["text_color"]

    st.write("---")
    st.subheader("🎉 完成したパンフレット")
    
    # ✨修正：1枚だけ生成した場合に巨大化しないよう、最低でも5分割のサイズに抑える
    num_cols = max(len(st.session_state.finished_pages), 5)
    cols = st.columns(num_cols)
    
    # ブラウザ上に完成した画像を並べて表示する
    for i, (page_data, page_img) in enumerate(st.session_state.finished_pages):
        with cols[i]:
            st.write(f"**{page_data['type']}**")
            st.image(page_img, use_container_width=True)

    # パワポファイルの作成準備と、サイズ（縦・横）の設定
    prs = Presentation()
    if orientation == "横向き (Landscape)":
        prs.slide_width, prs.slide_height = Inches(10), Inches(7.5)
    else:
        prs.slide_width, prs.slide_height = Inches(7.5), Inches(10)
    
# ✨ PowerPoint上に「自由に動かせる枠」を作るための便利機能（✨アイコン・ラベル・ゴールドフレーム・游明朝対応版）
        def add_placeholder_box(slide_obj, left, top, width, height, text, luxury_gold_rgb=None, show_icon=False, label_text="MAIN SLOT"):
            tx_box = slide_obj.shapes.add_textbox(left, top, width, height)
            
            tx_box.fill.solid()
            tx_box.fill.fore_color.rgb = RGBColor(245, 245, 245) # 予備の薄いグレー
            
            # ✨ 修正：ゴールドの金属質感フレームを追加
            if luxury_gold_rgb:
                tx_box.line.color.rgb = luxury_gold_rgb
            else:
                tx_box.line.color.rgb = RGBColor(185, 160, 110) # デフォルトゴールド
            tx_box.line.width = Pt(1)
            
            tf = tx_box.text_frame
            tf.word_wrap = True
            tf.clear()
            
            # プレースホルダーの用途ラベルとアイコンを追加
            if show_icon:
                p_label = tf.add_paragraph()
                p_label.text = label_text
                p_label.font.size = Pt(9)
                p_label.font.name = "游ゴシック"
                p_label.font.color.rgb = RGBColor(150, 150, 150) # 薄いグレー
                p_label.alignment = PP_ALIGN.CENTER
                # 小さなアイコン（ここではテキストで代用）
                p_icon = tf.add_paragraph()
                p_icon.text = "🔍" # アイコンの代わり
                p_icon.font.size = Pt(12)
                p_icon.alignment = PP_ALIGN.CENTER
            
                
            # ✨ 修正：鮮明で読みやすい書体（日本語は明朝体）に変更
            p_jp = tf.add_paragraph()
            clean_text = text.replace("【", "").replace("】", "").replace(" ", "")
            p_jp.text = clean_text
            p_jp.font.size = Pt(11) # 少し大きく
            p_jp.font.name = "游明朝" # 洗練された明朝体に変更
            p_jp.font.color.rgb = RGBColor(80, 80, 80) # ダークグレー
            p_jp.alignment = PP_ALIGN.CENTER

    # ✨ 駅徒歩用の円形グラフィック機能
    def add_station_info_circle(slide_obj, x, y, radius, text, color=(245, 240, 225)):
        circle = slide_obj.shapes.add_shape(MSO_SHAPE.OVAL, x, y, radius*2, radius*2)
        circle.fill.solid()
        circle.fill.fore_color.rgb = RGBColor(20, 20, 20)
        circle.line.color.rgb = RGBColor(color[0], color[1], color[2]) 
        circle.line.width = Pt(1.5)
        tf = circle.text_frame
        tf.word_wrap = True
        tf.clear()
        p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(color[0], color[1], color[2]) 
        p.font.name = "游明朝"
        p.alignment = PP_ALIGN.CENTER
        tf.margin_bottom = tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.05)

    # --- ✨修正：インデックス番号ではなく、page_dataのタイプでパワポを作成する ---
    for i, (page_data, page_img) in enumerate(st.session_state.finished_pages):
        p_type = page_data['type']
        slide = prs.slides.add_slide(prs.slide_layouts[6]) 
        
        img_io = BytesIO()
        page_img.save(img_io, format='PNG')
        img_io.seek(0)
        
        # aerial_map以外は背景画像を貼り付ける
        if p_type != 'aerial_map':
            slide.shapes.add_picture(img_io, 0, 0, width=prs.slide_width, height=prs.slide_height)

        # ────────── ① 表紙 ──────────
        if p_type == 'cover':
            prop_name = page_data.get('property_name_jp', '物件名') 
            address_text = page_data.get('city_town', '') 
            sub_copy = page_data.get('sub_copy', '') 
            station_info = page_data.get('station_info', '').replace('\\n', '\n')

            sw = prs.slide_width
            sh = prs.slide_height
            gold_color = RGBColor(215, 185, 140)
            luxury_font = "游明朝"

            top_bar_h = Inches(1.8) if orientation == "横向き (Landscape)" else Inches(2.2)
            top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, sw, top_bar_h)
            top_bar.fill.solid()
            top_bar.fill.fore_color.rgb = RGBColor(20, 20, 20)
            try: top_bar.fill.transparency = 0.4
            except: pass
            top_bar.line.fill.background() 

            bottom_bar_h = Inches(1.2) if orientation == "横向き (Landscape)" else Inches(1.5)
            bottom_bar_y = sh - bottom_bar_h
            bottom_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, bottom_bar_y, sw, bottom_bar_h)
            bottom_bar.fill.solid()
            bottom_bar.fill.fore_color.rgb = RGBColor(20, 20, 20)
            try: bottom_bar.fill.transparency = 0.4
            except: pass
            bottom_bar.line.fill.background()

            tx_addr = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(4.0), Inches(0.4))
            p_addr = tx_addr.text_frame.paragraphs[0]
            p_addr.text = address_text
            p_addr.font.size = Pt(13)
            p_addr.font.name = luxury_font
            p_addr.font.color.rgb = RGBColor(230, 230, 230)
            p_addr.alignment = PP_ALIGN.LEFT

            logo_y = Inches(0.6)
            if len(prop_name) > 15:
                logo_font_size = Pt(36) if orientation == "横向き (Landscape)" else Pt(32)
            elif len(prop_name) > 10:
                logo_font_size = Pt(48) if orientation == "横向き (Landscape)" else Pt(42)
            else:
                logo_font_size = Pt(60) if orientation == "横向き (Landscape)" else Pt(54)

            tx_logo = slide.shapes.add_textbox(Inches(0.35), logo_y, sw - Inches(0.8), Inches(2.0))
            tf_logo = tx_logo.text_frame
            tf_logo.word_wrap = False
            p_logo = tf_logo.paragraphs[0]
            p_logo.text = prop_name
            p_logo.font.size = logo_font_size
            p_logo.font.bold = True
            p_logo.font.name = luxury_font
            p_logo.font.color.rgb = gold_color
            p_logo.alignment = PP_ALIGN.LEFT

            sub_y = bottom_bar_y + Inches(0.3)
            tx_sub = slide.shapes.add_textbox(Inches(1.5), sub_y, sw - Inches(1.8), Inches(1.0))
            tf_sub = tx_sub.text_frame
            tf_sub.word_wrap = True
            p_sub = tf_sub.paragraphs[0]
            p_sub.text = sub_copy
            p_sub.font.size = Pt(13)
            p_sub.font.name = luxury_font
            p_sub.font.color.rgb = RGBColor(240, 240, 240)
            p_sub.alignment = PP_ALIGN.CENTER

            circle_y = bottom_bar_y - Inches(0.5)
            add_station_info_circle(slide, Inches(0.4), circle_y, Inches(0.55), station_info, color=(215, 185, 140))

        # ────────── ② コンセプト ──────────
        elif p_type == 'aerial_map':
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(245, 240, 225)
            
            headline = "THE CONCEPT"
            sub_headline = page_data.get('sub_headline', '五感を満たす、静謐の邸宅。')
            main_text = page_data.get('main_text', '')
            bronze_rgb = RGBColor(215, 185, 140)
            headline_font = "Arial"
            main_font = "游ゴシック"
            
            if orientation == "横向き (Landscape)":
                tx_box_head = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(4.5), Inches(0.6))
                line1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.3), Inches(4.5), Pt(1))
                line1.fill.solid()
                line1.fill.fore_color.rgb = bronze_rgb
                line1.line.color.rgb = bronze_rgb
                tx_box_sub = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(4.5), Inches(0.6))
                tx_box_main = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(4.0), Inches(3.0))
                
                if "concept_img" in st.session_state and st.session_state.concept_img:
                    img_io_concept = BytesIO()
                    st.session_state.concept_img.save(img_io_concept, format='PNG')
                    img_io_concept.seek(0)
                    slide.shapes.add_picture(img_io_concept, Inches(4.5), Inches(0), width=Inches(5.5), height=Inches(7.5))
                else:
                    add_placeholder_box(slide, Inches(4.5), Inches(0), Inches(5.5), Inches(7.5), "メイン画像")
                add_placeholder_box(slide, Inches(3.0), Inches(4.5), Inches(2.5), Inches(2.0), "写真1")
                add_placeholder_box(slide, Inches(7.5), Inches(0.5), Inches(2.0), Inches(2.5), "写真2")
            else:
                tx_box_head = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(6.5), Inches(0.6))
                line1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.3), Inches(6.5), Pt(1))
                line1.fill.solid()
                line1.fill.fore_color.rgb = bronze_rgb
                line1.line.color.rgb = bronze_rgb
                tx_box_sub = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(6.5), Inches(0.6))
                tx_box_main = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(6.5), Inches(1.5))
                
                if "concept_img" in st.session_state and st.session_state.concept_img:
                    img_io_concept = BytesIO()
                    st.session_state.concept_img.save(img_io_concept, format='PNG')
                    img_io_concept.seek(0)
                    slide.shapes.add_picture(img_io_concept, Inches(0), Inches(4.5), width=Inches(7.5), height=Inches(5.5))
                else:
                    add_placeholder_box(slide, Inches(0), Inches(4.5), Inches(7.5), Inches(5.5), "メイン画像")
                add_placeholder_box(slide, Inches(4.5), Inches(3.5), Inches(2.5), Inches(2.0), "写真1")
                add_placeholder_box(slide, Inches(0.5), Inches(7.5), Inches(2.5), Inches(2.0), "写真2")

            tf_head = tx_box_head.text_frame
            p_head = tf_head.paragraphs[0]
            p_head.text = headline
            p_head.font.name = headline_font
            p_head.font.size = Pt(36)
            p_head.font.bold = True
            p_head.font.color.rgb = RGBColor(0,0,0)

            tf_sub = tx_box_sub.text_frame
            p_sub = tf_sub.paragraphs[0]
            p_sub.text = sub_headline
            p_sub.font.name = main_font
            p_sub.font.size = Pt(24)
            p_sub.font.bold = True
            p_sub.font.color.rgb = RGBColor(0,0,0)

            tf_main = tx_box_main.text_frame
            tf_main.word_wrap = True
            p_main = tf_main.paragraphs[0]
            p_main.text = main_text
            p_main.font.name = main_font
            p_main.font.size = Pt(13)
            p_main.font.color.rgb = RGBColor(80,80,80)

        # ────────── ③ MAP & ACCESS ──────────
        elif p_type == 'access':
            acc_text = page_data.get('access_info')
            if not acc_text: acc_text = "交通情報がありません"
            elif isinstance(acc_text, list): acc_text = "\n".join(str(x) for x in acc_text)
            else: acc_text = str(acc_text)

            life_text = page_data.get('life_info')
            if not life_text: life_text = "周辺施設情報がありません"
            elif isinstance(life_text, list): life_text = "\n".join(str(x) for x in life_text)
            else: life_text = str(life_text)

            gold_color = RGBColor(180, 150, 80)
            navy_color = RGBColor(20, 30, 60)
            text_color_rgb = RGBColor(40, 40, 40)
            
            gen_map_io = None
            if map_file:
                from PIL import ImageOps, ImageEnhance
                try:
                    map_img = Image.open(map_file).convert("RGB")
                    gray_img = map_img.convert("L")
                    inverted_img = ImageOps.invert(gray_img)
                    navy_bg = (15, 25, 45)
                    gold_line = (215, 185, 140)
                    styled_map = ImageOps.colorize(inverted_img, black=navy_bg, white=gold_line)
                    enhancer = ImageEnhance.Contrast(styled_map)
                    final_map = enhancer.enhance(1.3)
                    gen_map_io = BytesIO()
                    final_map.save(gen_map_io, format='PNG')
                    gen_map_io.seek(0)
                except Exception as e:
                    pass

            if orientation == "横向き (Landscape)":
                tx_box_head = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(5.0), Inches(0.8))
                tx_box_sub = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(5.0), Inches(0.4))
                line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(9.0), Pt(1.5))
                line.fill.solid()
                line.fill.fore_color.rgb = gold_color
                line.line.color.rgb = gold_color

                if gen_map_io:
                    slide.shapes.add_picture(gen_map_io, Inches(0.5), Inches(1.7), width=Inches(9.0), height=Inches(4.2))
                else:
                    add_placeholder_box(slide, Inches(0.5), Inches(1.7), Inches(9.0), Inches(4.2), "【 MAP画像 挿入枠 】\n※画像をアップロードすると自動で高級な地図が生成されます")

                tx_box_acc = slide.shapes.add_textbox(Inches(0.5), Inches(6.1), Inches(4.3), Inches(1.2))
                tx_box_life = slide.shapes.add_textbox(Inches(5.2), Inches(6.1), Inches(4.3), Inches(1.2))
            else:
                tx_box_head = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(6.5), Inches(0.8))
                tx_box_sub = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(6.5), Inches(0.4))
                line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.6), Inches(6.5), Pt(1.5))
                line.fill.solid()
                line.fill.fore_color.rgb = gold_color
                line.line.color.rgb = gold_color

                if gen_map_io:
                    slide.shapes.add_picture(gen_map_io, Inches(0.5), Inches(1.8), width=Inches(6.5), height=Inches(5.0))
                else:
                    add_placeholder_box(slide, Inches(0.5), Inches(1.8), Inches(6.5), Inches(5.0), "【 MAP画像 挿入枠 】\n※画像をアップロードすると自動で高級な地図が生成されます")
                
                tx_box_acc = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(3.1), Inches(2.5))
                tx_box_life = slide.shapes.add_textbox(Inches(3.9), Inches(7.0), Inches(3.1), Inches(2.5))

            tf_head = tx_box_head.text_frame
            p_head = tf_head.paragraphs[0]
            p_head.text = "MAP & ACCESS"
            p_head.font.size = Pt(40)
            p_head.font.bold = True
            p_head.font.color.rgb = navy_color
            p_head.font.name = "Arial"

            tf_sub = tx_box_sub.text_frame
            p_sub = tf_sub.paragraphs[0]
            p_sub.text = "周辺環境・アクセス"
            p_sub.font.size = Pt(20)
            p_sub.font.bold = True
            p_sub.font.color.rgb = gold_color
            p_sub.font.name = "游ゴシック"

            tf_acc = tx_box_acc.text_frame
            tf_acc.word_wrap = True
            tf_acc.clear() 
            p_acc_head = tf_acc.add_paragraph()
            p_acc_head.text = "■ 交通アクセス"
            p_acc_head.font.size = Pt(16)
            p_acc_head.font.bold = True
            p_acc_head.font.color.rgb = gold_color 
            p_acc_head.font.name = "游ゴシック"
            p_acc_body = tf_acc.add_paragraph()
            p_acc_body.text = acc_text
            p_acc_body.font.size = Pt(13)
            p_acc_body.font.color.rgb = text_color_rgb 
            p_acc_body.font.name = "游ゴシック"

            tf_life = tx_box_life.text_frame
            tf_life.word_wrap = True
            tf_life.clear() 
            p_life_head = tf_life.add_paragraph()
            p_life_head.text = "■ Life Information"
            p_life_head.font.size = Pt(16)
            p_life_head.font.bold = True
            p_life_head.font.color.rgb = gold_color 
            p_life_head.font.name = "Arial"
            p_life_body = tf_life.add_paragraph()
            p_life_body.text = life_text
            p_life_body.font.size = Pt(13)
            p_life_body.font.color.rgb = text_color_rgb
            p_life_body.font.name = "游ゴシック"

# ────────── ④ 間取り図 ──────────
        elif p_type == 'floor_plan':
            prop_name_jp = page_data.get('property_name_jp', '物件名')
            land_area = page_data.get('land_area', '')
            building_area = page_data.get('building_area', '')
            price_jp = page_data.get('price_jp', '--- 万円')

            # ✨ 修正：お手本に合わせた高級感のあるブロンズゴールドと游明朝の指定
            luxury_gold = RGBColor(185, 160, 110)
            
            # ✨ 修正1：テキストボックスの座標と幅を調整し、左右のエリアを明確に分ける
            if orientation == "横向き (Landscape)":
                tb_title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(5.8), Inches(1.2))
                tb_area = slide.shapes.add_textbox(Inches(6.5), Inches(0.1), Inches(3.0), Inches(0.8))
                tb_price = slide.shapes.add_textbox(Inches(6.5), Inches(0.8), Inches(3.0), Inches(0.8))
            else:
                tb_title = slide.shapes.add_textbox(Inches(0.4), Inches(0.3), Inches(4.2), Inches(1.2))
                tb_area = slide.shapes.add_textbox(Inches(4.6), Inches(0.1), Inches(2.7), Inches(0.8))
                tb_price = slide.shapes.add_textbox(Inches(4.6), Inches(0.8), Inches(2.7), Inches(0.8))
             
            # 物件タイトル
            tf_title = tb_title.text_frame
            tf_title.word_wrap = True # 長すぎる場合は改行して被りを防ぐ
            p_title = tf_title.paragraphs[0]
            p_title.text = prop_name_jp
            
            # ✨ 修正：号室を含んで長くなっても見切れないよう、文字数に合わせてサイズを細かく調整
            name_len = len(prop_name_jp)
            if name_len > 20:
                p_title.font.size = Pt(20)
            elif name_len > 14:
                p_title.font.size = Pt(24)
            elif name_len > 10:
                p_title.font.size = Pt(28)
            else:
                p_title.font.size = Pt(32)
                
            p_title.font.name = "游明朝"
            p_title.font.bold = True
            p_title.font.color.rgb = luxury_gold
            try: p_title.font.shadow = True # 文字に影を付ける
            except: pass
            
            # 土地・建物面積
            tf_area = tb_area.text_frame
            tf_area.clear()
            p_land = tf_area.add_paragraph()
            p_land.text = land_area
            p_land.font.size = Pt(16) # ✨ 修正4：面積の文字を Pt(11) から Pt(16) に拡大
            p_land.font.color.rgb = luxury_gold
            p_land.alignment = PP_ALIGN.RIGHT
            try: p_land.font.shadow = True 
            except: pass
            
            if building_area:
                p_build = tf_area.add_paragraph()
                p_build.text = building_area
                p_build.font.size = Pt(16) # ✨ 修正4：面積の文字を Pt(11) から Pt(16) に拡大
                p_build.font.color.rgb = luxury_gold
                p_build.alignment = PP_ALIGN.RIGHT
                try: p_build.font.shadow = True 
                except: pass
            
            # 価格
            tf_price = tb_price.text_frame
            p_price = tf_price.paragraphs[0]
            p_price.text = price_jp
            p_price.font.size = Pt(36)
            p_price.font.name = "游明朝"
            p_price.font.bold = True
            p_price.font.color.rgb = luxury_gold
            p_price.alignment = PP_ALIGN.RIGHT
            try: p_price.font.shadow = True # 価格にも影を付ける
            except: pass

# ────────── ⑥ 内観ギャラリー（✨P.7：金色の線あり・5枚枠） ──────────
        elif p_type == 'interior':
            sw, sh = prs.slide_width, prs.slide_height
            luxury_gold = RGBColor(185, 160, 110)
            
            # 座標設定
            if orientation == "横向き (Landscape)":
                en_y, jp_y, bar_y, main_y, main_h = Inches(0.2), Inches(0.7), Inches(1.4), Inches(1.6), Inches(2.5)
                sub1_y, sub2_y, sub_h = Inches(4.3), Inches(5.8), Inches(1.3)
            else:
                en_y, jp_y, bar_y, main_y, main_h = Inches(0.3), Inches(0.8), Inches(1.5), Inches(1.7), Inches(3.7)
                sub1_y, sub2_y, sub_h = Inches(5.6), Inches(7.6), Inches(1.8)

            # ゴールドバー（P.7は必要）
            gold_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, bar_y, sw, Pt(6))
            gold_bar.fill.solid()
            gold_bar.fill.fore_color.rgb = luxury_gold
            gold_bar.line.fill.background()

            # 5枚の画像枠
            if orientation == "横向き (Landscape)":
                add_placeholder_box(slide, Inches(0.5), main_y, Inches(9.0), main_h, "【 メイン画像 】", luxury_gold, show_icon=True, label_text="MAIN SLOT")
                add_placeholder_box(slide, Inches(0.5), sub1_y, Inches(4.3), sub_h, "【 サブ画像1 】", luxury_gold, show_icon=True, label_text="SUB 1")
                add_placeholder_box(slide, Inches(5.2), sub1_y, Inches(4.3), sub_h, "【 サブ画像2 】", luxury_gold, show_icon=True, label_text="SUB 2")
                add_placeholder_box(slide, Inches(0.5), sub2_y, Inches(4.3), sub_h, "【 サブ画像3 】", luxury_gold, show_icon=True, label_text="SUB 3")
                add_placeholder_box(slide, Inches(5.2), sub2_y, Inches(4.3), sub_h, "【 サブ画像4 】", luxury_gold, show_icon=True, label_text="SUB 4")
            else:
                add_placeholder_box(slide, Inches(0.5), main_y, Inches(6.5), main_h, "【 メイン画像 】", luxury_gold, show_icon=True, label_text="MAIN SLOT")
                add_placeholder_box(slide, Inches(0.5), sub1_y, Inches(3.1), sub_h, "【 サブ1 】", luxury_gold, show_icon=True, label_text="SUB 1")
                add_placeholder_box(slide, Inches(3.9), sub1_y, Inches(3.1), sub_h, "【 サブ2 】", luxury_gold, show_icon=True, label_text="SUB 2")
                add_placeholder_box(slide, Inches(0.5), sub2_y, Inches(3.1), sub_h, "【 サブ3 】", luxury_gold, show_icon=True, label_text="SUB 3")
                add_placeholder_box(slide, Inches(3.9), sub2_y, Inches(3.1), sub_h, "【 サブ4 】", luxury_gold, show_icon=True, label_text="SUB 4")

            # ヘッダー（P.7は必要）
            tx_head_en = slide.shapes.add_textbox(sw * 0.05, en_y, sw * 0.5, Inches(0.5))
            tx_head_en.text_frame.paragraphs[0].text = "INTERIOR GALLERY"
            tx_head_en.text_frame.paragraphs[0].font.size, tx_head_en.text_frame.paragraphs[0].font.color.rgb = Pt(28), RGBColor(255, 255, 255)
            
            tx_head_jp = slide.shapes.add_textbox(sw * 0.05, jp_y, sw * 0.4, Inches(0.6))
            tx_head_jp.text_frame.paragraphs[0].text = "内観ギャラリー"
            tx_head_jp.text_frame.paragraphs[0].font.size, tx_head_jp.text_frame.paragraphs[0].font.name, tx_head_jp.text_frame.paragraphs[0].font.color.rgb = Pt(16), "游明朝", luxury_gold

        # ────────── ⑦ 会社案内（✨P.8：文字・線を消して3枚大型化版） ──────────
        elif p_type == 'company':
            sw, sh = prs.slide_width, prs.slide_height
            luxury_gold = RGBColor(185, 160, 110)
            gold_main = luxury_gold # NameError回避

            # ヘッダーを削除したため、上から大きく配置
            if orientation == "横向き (Landscape)":
                main_y, main_h = Inches(0.4), Inches(4.2)
                sub_y, sub_h = Inches(4.8), Inches(2.1)
                sub_w = Inches(4.3)
            else:
                main_y, main_h = Inches(0.5), Inches(5.5)
                sub_y, sub_h = Inches(6.2), Inches(3.0)
                sub_w = Inches(3.1)

            # 3つの大型画像枠
            if orientation == "横向き (Landscape)":
                add_placeholder_box(slide, Inches(0.5), main_y, Inches(9.0), main_h, "【 会社案内 メイン画像 】", luxury_gold, show_icon=True, label_text="COMPANY MAIN")
                add_placeholder_box(slide, Inches(0.5), sub_y, sub_w, sub_h, "【 オフィス写真 】", luxury_gold, show_icon=True, label_text="OFFICE")
                add_placeholder_box(slide, Inches(5.2), sub_y, sub_w, sub_h, "【 スタッフ写真 】", luxury_gold, show_icon=True, label_text="STAFF")
            else:
                add_placeholder_box(slide, Inches(0.5), main_y, Inches(6.5), main_h, "【 会社案内 メイン画像 】", luxury_gold, show_icon=True, label_text="COMPANY MAIN")
                add_placeholder_box(slide, Inches(0.5), sub_y, sub_w, sub_h, "【 オフィス写真 】", luxury_gold, show_icon=True, label_text="OFFICE")
                add_placeholder_box(slide, Inches(3.9), sub_y, sub_w, sub_h, "【 スタッフ写真 】", luxury_gold, show_icon=True, label_text="STAFF")

            # ブランドメッセージ（最下部に配置）
            tx_slogan = slide.shapes.add_textbox(0, sh * 0.94, sw, Inches(0.4))
            p_slogan = tx_slogan.text_frame.paragraphs[0]
            p_slogan.text = "「住まい」を通じて、お客様の人生に確かな価値を。"
            p_slogan.font.size, p_slogan.font.name = Pt(14), "游明朝"
            p_slogan.font.color.rgb = RGBColor(245, 240, 225)
            p_slogan.alignment = PP_ALIGN.CENTER

# --- スライド2：Contact & Thank you（THANK YOUページ） ---
            slide_thanks = prs.slides.add_slide(prs.slide_layouts[6]) 
            fill_thanks = slide_thanks.background.fill
            fill_thanks.solid()
            fill_thanks.fore_color.rgb = RGBColor(252, 248, 242)

            top_line = slide_thanks.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, sw, Pt(6))
            top_line.fill.solid()
            top_line.fill.fore_color.rgb = gold_main
            top_line.line.fill.background()

            tx_ty = slide_thanks.shapes.add_textbox(0, sh * 0.12, sw, Inches(1.0)) # 少し上に移動
            p_ty = tx_ty.text_frame.paragraphs[0]
            p_ty.text = "THANK YOU"
            p_ty.font.size, p_ty.font.name, p_ty.alignment = Pt(48), "Times New Roman", PP_ALIGN.CENTER

            # ✨ 修正1：店舗案内図の枠を少し小さくし、位置を上に上げて重なりを防止
            map_box_w, map_box_h = Inches(5.5), Inches(2.2) 
            add_placeholder_box(slide_thanks, (sw - map_box_w) / 2, sh * 0.30, map_box_w, map_box_h, "店舗案内図（地図）")

            # ✨ 修正2：フリーダイヤル枠の開始位置（box_y）を sh * 0.68 まで下げて重なりを解消
            box_w, box_h, box_y = Inches(6.5), Inches(1.8), sh * 0.68 
            box = slide_thanks.shapes.add_shape(MSO_SHAPE.RECTANGLE, (sw - box_w) / 2, box_y, box_w, box_h)
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(255, 255, 255)
            box.line.color.rgb = RGBColor(240, 240, 240)
            box.line.width = Pt(1)

            # テキストボックス
            tx_contact = slide_thanks.shapes.add_textbox((sw - box_w) / 2, box_y, box_w, box_h)
            tf_contact = tx_contact.text_frame
            tf_contact.vertical_anchor = MSO_ANCHOR.MIDDLE # 上下中央揃え 
            
            # ✨ 修正3：テキスト枠自体の余白をゼロにして、上下中央をより正確にする
            tf_contact.margin_top = tf_contact.margin_bottom = 0
            tf_contact.clear()
            
            p_tel = tf_contact.add_paragraph()
            p_tel.text = f"フリーダイヤル：{page_data.get('tel', '')}"
            p_tel.font.size, p_tel.font.bold, p_tel.font.color.rgb, p_tel.alignment = Pt(26), True, gold_main, PP_ALIGN.CENTER

            p_info = tf_contact.add_paragraph()
            p_info.text = f"{page_data.get('license', '')}  |  {page_data.get('address', '')}"
            p_info.font.size, p_info.font.color.rgb, p_info.alignment = Pt(10), RGBColor(120, 120, 120), PP_ALIGN.CENTER

    # ✨ 修正：完成後のダウンロード処理
    pptx_out = BytesIO()
    prs.save(pptx_out)
    pptx_out.seek(0)

    st.write("---")
    st.download_button(
        label="📥 完成したパンフレット（PowerPoint）を保存する",
        data=pptx_out,
        file_name=f"パンフレット_{st.session_state.current_file.split('.')[0]}.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )